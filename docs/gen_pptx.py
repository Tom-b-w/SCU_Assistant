"""
SCU Assistant 开题报告 PPT — 腾讯奖学金模板风格
蓝色渐变 + 白色内容区 + 图表丰富 + 数据支撑
最终输出 .pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData, ChartData
import os

# ============================================================
# 品牌色系 — 腾讯蓝 + 科技感
# ============================================================
TENC_DARK   = RGBColor(0x00, 0x27, 0x66)
TENC_BLUE   = RGBColor(0x00, 0x52, 0xD9)
TENC_LIGHT  = RGBColor(0x3B, 0x82, 0xF6)
TENC_BG     = RGBColor(0xEB, 0xF2, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF8, 0xFA, 0xFF)
LIGHT_GRAY  = RGBColor(0xE2, 0xE8, 0xF0)
MID_GRAY    = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT    = RGBColor(0x1E, 0x29, 0x3B)
SUB_TEXT    = RGBColor(0x64, 0x74, 0x8B)
SUCCESS     = RGBColor(0x00, 0xB5, 0x78)
WARN        = RGBColor(0xFF, 0x88, 0x00)
DANGER      = RGBColor(0xE3, 0x4D, 0x59)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
GOLD        = RGBColor(0xD9, 0x77, 0x06)
SCU_RED     = RGBColor(0xC4, 0x12, 0x30)
CYAN        = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ============================================================
# 工具函数
# ============================================================

def _add_rect(slide, left, top, w, h, fill=None, border=None, border_w=Pt(0.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    return s

def _add_rounded(slide, left, top, w, h, fill=None, border=None, border_w=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    return s

def _add_circle(slide, left, top, size, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def _tb(slide, left, top, w, h, text, sz=16, color=DARK_TEXT, bold=False,
        align=PP_ALIGN.LEFT, font="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    r = tf.paragraphs[0].add_run()
    r.text = text; r.font.size = Pt(sz); r.font.color.rgb = color
    r.font.bold = bold; r.font.name = font
    try:
        tf.paragraphs[0].font.name = font
    except:
        pass
    return tb

def _add_para(tf, text, sz=14, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT,
              space_before=Pt(2), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.alignment = align; p.space_before = space_before; p.space_after = space_after
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = "Microsoft YaHei"
    return p

def _page_num(slide, num, total):
    _tb(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
        f"{num}/{total}", sz=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def _blue_header(slide):
    """蓝色顶栏"""
    _add_rect(slide, 0, 0, SW, Inches(1.05), fill=TENC_BLUE)

def _section_title(slide, num, title, subtitle=""):
    _blue_header(slide)
    _tb(slide, Inches(0.7), Inches(0.2), Inches(1), Inches(0.6),
        num, sz=32, color=WHITE, bold=True, font="Segoe UI")
    _tb(slide, Inches(1.6), Inches(0.25), Inches(8), Inches(0.55),
        title, sz=24, color=WHITE, bold=True)
    if subtitle:
        _tb(slide, Inches(1.6), Inches(0.65), Inches(8), Inches(0.3),
            subtitle, sz=11, color=RGBColor(0xBB,0xD5,0xFF))

def _white_card(slide, left, top, w, h, border=LIGHT_GRAY):
    return _add_rounded(slide, left, top, w, h, fill=WHITE, border=border)

def _color_top_card(slide, left, top, w, h, accent_color, title, items, icon=""):
    """带色彩顶栏的卡片"""
    _add_rounded(slide, left, top, w, h, fill=WHITE, border=LIGHT_GRAY)
    _add_rect(slide, left+Inches(0.02), top+Inches(0.02), w-Inches(0.04), Inches(0.38),
              fill=accent_color)
    _tb(slide, left+Inches(0.12), top+Inches(0.05), w-Inches(0.2), Inches(0.32),
        f"{icon}  {title}" if icon else title, sz=11, color=WHITE, bold=True)
    y = top + Inches(0.5)
    for item in items:
        _tb(slide, left+Inches(0.12), y, w-Inches(0.24), Inches(0.28),
            f"• {item}", sz=9, color=SUB_TEXT)
        y += Inches(0.28)

def _stat_card(slide, left, top, w, h, number, label, num_color=TENC_BLUE):
    _white_card(slide, left, top, w, h)
    _tb(slide, left, top+Inches(0.15), w, Inches(0.5),
        number, sz=30, color=num_color, bold=True, align=PP_ALIGN.CENTER, font="Segoe UI")
    _tb(slide, left, top+Inches(0.65), w, Inches(0.3),
        label, sz=10, color=SUB_TEXT, align=PP_ALIGN.CENTER)

def _progress_bar(slide, left, top, total_w, h, ratio, color, bg=LIGHT_GRAY):
    _add_rounded(slide, left, top, total_w, h, fill=bg)
    if ratio > 0:
        _add_rounded(slide, left, top, int(total_w * ratio), h, fill=color)

TOTAL = 14

# ============================================================
# Slide 1: 封面
# ============================================================
def s01_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 深蓝背景
    _add_rect(slide, 0, 0, SW, SH, fill=TENC_DARK)
    # 左侧渐变色块
    # 用矩形模拟渐变色块
    _add_rect(slide, 0, 0, Inches(6.5), SH, fill=TENC_BLUE)
    # 右上装饰圆 (浅色低透明度) — 简化处理，不设透明度
    from pptx.oxml.ns import qn
    c = _add_circle(slide, Inches(10), Inches(-0.5), Inches(4), RGBColor(0x1A, 0x3D, 0x7A))

    # 右下装饰
    c2 = _add_circle(slide, Inches(8.5), Inches(4.5), Inches(5), RGBColor(0x10, 0x30, 0x60))

    # 底部线
    _add_rect(slide, 0, Inches(7.3), SW, Inches(0.03), fill=RGBColor(0xFF,0xFF,0xFF))

    # 文字
    _tb(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
        "SCU Assistant", sz=48, color=WHITE, bold=True, font="Segoe UI")
    _tb(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.7),
        "四川大学智能校园助手", sz=30, color=RGBColor(0xBB,0xD5,0xFF), bold=True)

    # 分割线
    _add_rect(slide, Inches(1.5), Inches(3.8), Inches(2.5), Inches(0.03), fill=TENC_LIGHT)

    _tb(slide, Inches(1.5), Inches(4.1), Inches(8), Inches(0.4),
        "项目开题报告", sz=18, color=RGBColor(0xDD,0xDD,0xFF))
    _tb(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.4),
        "四川大学计算机学院  |  软件工程课程设计  |  2026 春季学期",
        sz=12, color=MID_GRAY)

    # SCU 标志圆
    c3 = _add_circle(slide, Inches(10.5), Inches(2.5), Inches(1.6), SCU_RED)
    c4 = _add_circle(slide, Inches(10.7), Inches(2.7), Inches(1.2), TENC_DARK)
    _tb(slide, Inches(10.7), Inches(2.95), Inches(1.2), Inches(0.7),
        "SCU\nAI+", sz=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font="Segoe UI")

    _page_num(slide, 1, TOTAL)

# ============================================================
# Slide 2: 目录
# ============================================================
def s02_toc():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "", "目录", "CONTENTS")
    # 重绘顶栏
    _add_rect(slide, 0, 0, SW, Inches(0.9), fill=TENC_BLUE)
    _tb(slide, Inches(0.8), Inches(0.2), Inches(6), Inches(0.5),
        "目录  CONTENTS", sz=22, color=WHITE, bold=True)

    items_l = [
        ("01", "项目背景与痛点", "校园信息碎片化现状分析"),
        ("02", "产品定位与目标", "AI驱动的一站式校园助手"),
        ("03", "核心功能模块", "6大核心模块详解"),
        ("04", "系统架构设计", "三层架构 + 容器化部署"),
        ("05", "技术选型", "前端/后端/AI/DevOps 技术栈"),
    ]
    items_r = [
        ("06", "团队成员与分工", "8人4组敏捷协作"),
        ("07", "开发计划与里程碑", "16周敏捷迭代路线图"),
        ("08", "技术难点与风险", "风险识别与应对策略"),
        ("09", "当前进展", "P0核心模块完成度"),
        ("10", "总结与展望", "项目亮点与下一步计划"),
    ]

    for col_items, start_x in [(items_l, Inches(0.8)), (items_r, Inches(6.9))]:
        for i, (num, title, desc) in enumerate(col_items):
            y = Inches(1.3) + Inches(1.1) * i
            _white_card(slide, start_x, y, Inches(5.5), Inches(0.9))
            # 蓝色数字
            _add_rounded(slide, start_x + Inches(0.15), y + Inches(0.15),
                         Inches(0.6), Inches(0.6), fill=TENC_BLUE)
            _tb(slide, start_x + Inches(0.15), y + Inches(0.2),
                Inches(0.6), Inches(0.5), num, sz=20, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER, font="Segoe UI")
            _tb(slide, start_x + Inches(0.9), y + Inches(0.12),
                Inches(4), Inches(0.35), title, sz=15, color=DARK_TEXT, bold=True)
            _tb(slide, start_x + Inches(0.9), y + Inches(0.48),
                Inches(4), Inches(0.3), desc, sz=10, color=SUB_TEXT)

    _page_num(slide, 2, TOTAL)

# ============================================================
# Slide 3: 项目背景
# ============================================================
def s03_background():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "01", "项目背景与痛点", "校园信息化现状分析")

    # 左侧：痛点描述
    _tb(slide, Inches(0.8), Inches(1.2), Inches(6), Inches(0.4),
        "四川大学在校学生超过 5 万人，日常需频繁访问多个独立平台获取校园信息。",
        sz=13, color=SUB_TEXT)

    pain_points = [
        ("信息碎片化", "课表、校车、食堂分布在 5+ 个平台", DANGER, "!"),
        ("系统体验差", "教务 UI 落后，移动端适配极差", WARN, "✕"),
        ("缺乏智能化", "无个性推荐，无统一问答入口", PURPLE, "?"),
        ("数据不互通", "各系统数据孤岛，时效性差", TENC_LIGHT, "#"),
    ]

    card_w = Inches(2.8)
    start_x = Inches(0.6)
    gap = Inches(0.25)
    cy = Inches(1.8)

    for i, (title, desc, color, icon) in enumerate(pain_points):
        x = start_x + (card_w + gap) * i
        _white_card(slide, x, cy, card_w, Inches(2.2))
        # 图标圆
        c = _add_circle(slide, x + Inches(1.05), cy + Inches(0.2), Inches(0.6), color)
        _tb(slide, x + Inches(1.05), cy + Inches(0.25), Inches(0.6), Inches(0.5),
            icon, sz=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Segoe UI")
        _tb(slide, x + Inches(0.15), cy + Inches(0.9), card_w - Inches(0.3), Inches(0.35),
            title, sz=14, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, x + Inches(0.15), cy + Inches(1.3), card_w - Inches(0.3), Inches(0.7),
            desc, sz=10, color=SUB_TEXT, align=PP_ALIGN.CENTER)

    # 数据支撑卡片
    _tb(slide, Inches(0.8), Inches(4.25), Inches(4), Inches(0.35),
        "数据支撑", sz=16, color=TENC_BLUE, bold=True)

    stats = [
        ("50,000+", "在校学生总数", TENC_BLUE),
        ("5+", "日均使用平台数", WARN),
        ("73%", "反馈查课不方便", DANGER),
        ("89%", "希望有智能助手", SUCCESS),
        ("0", "现有同类产品", GOLD),
    ]
    sx = Inches(0.6)
    sw = Inches(2.3)
    sg = Inches(0.18)
    for i, (num, label, clr) in enumerate(stats):
        x = sx + (sw + sg) * i
        _stat_card(slide, x, Inches(4.65), sw, Inches(1.0), num, label, clr)

    # 右侧结论
    _white_card(slide, Inches(0.6), Inches(5.9), Inches(11.9), Inches(1.1), border=TENC_BLUE)
    _add_rect(slide, Inches(0.62), Inches(5.92), Inches(0.08), Inches(1.06), fill=TENC_BLUE)
    _tb(slide, Inches(1.0), Inches(6.05), Inches(10), Inches(0.4),
        "结论：市场空白，川大尚无一款集成式智能校园助手产品",
        sz=15, color=TENC_BLUE, bold=True)
    _tb(slide, Inches(1.0), Inches(6.4), Inches(10), Inches(0.35),
        "学生迫切需要一个统一入口，通过 AI 高效获取课表、食堂、校车、DDL 等校园信息，减少平台切换、提升获取效率",
        sz=10, color=SUB_TEXT)

    _page_num(slide, 3, TOTAL)

# ============================================================
# Slide 4: 产品定位
# ============================================================
def s04_positioning():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "02", "产品定位与核心目标")

    # 左：核心定位
    _white_card(slide, Inches(0.6), Inches(1.3), Inches(6.0), Inches(2.5), border=TENC_BLUE)
    _add_rect(slide, Inches(0.62), Inches(1.32), Inches(0.06), Inches(2.46), fill=TENC_BLUE)
    _tb(slide, Inches(0.9), Inches(1.4), Inches(2), Inches(0.3),
        "核心定位", sz=12, color=TENC_BLUE, bold=True)
    _tb(slide, Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.5),
        "AI 驱动的一站式智能校园助手", sz=20, color=DARK_TEXT, bold=True)
    _tb(slide, Inches(0.9), Inches(2.35), Inches(5.3), Inches(1.2),
        "• 自然语言对话 + 传统页面导航 双模式交互\n• 对接真实教务系统数据，非 Mock 演示\n• 个性化记忆系统，越用越懂你",
        sz=12, color=SUB_TEXT)

    # 右：目标用户
    _white_card(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(2.5))
    _tb(slide, Inches(7.2), Inches(1.4), Inches(3), Inches(0.3),
        "目标用户", sz=12, color=TENC_BLUE, bold=True)
    users = [
        "四川大学全日制本科生 & 研究生 (50,000+)",
        "日均使用教务/食堂/校车等校园服务的学生",
        "期望通过 AI 提升信息获取效率的用户",
    ]
    for i, u in enumerate(users):
        y = Inches(1.85) + Inches(0.5) * i
        _add_circle(slide, Inches(7.2), y + Inches(0.04), Inches(0.18), SUCCESS)
        _tb(slide, Inches(7.55), y, Inches(4.8), Inches(0.35), u, sz=12, color=SUB_TEXT)

    # 四个核心目标
    goals = [
        ("效率提升", "一句话查课表/成绩\n减少 80% 操作步骤", TENC_BLUE, "⚡"),
        ("智能推荐", "AI 意图识别\n个性化食堂/选课推荐", PURPLE, "🧠"),
        ("数据整合", "教务/食堂/校车/DDL\n一站式信息聚合", GOLD, "📊"),
        ("体验升级", "现代化 UI 设计\n深色主题 + 移动优先", SUCCESS, "🎨"),
    ]
    gw = Inches(2.85)
    gg = Inches(0.23)
    gy = Inches(4.15)
    for i, (title, desc, color, icon) in enumerate(goals):
        gx = Inches(0.6) + (gw + gg) * i
        _white_card(slide, gx, gy, gw, Inches(2.85))
        # 顶色条
        _add_rect(slide, gx + Inches(0.02), gy + Inches(0.02),
                  gw - Inches(0.04), Inches(0.06), fill=color)
        _tb(slide, gx + Inches(0.15), gy + Inches(0.25), gw - Inches(0.3), Inches(0.35),
            f"{icon}  {title}", sz=15, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
        # 分割线
        _add_rect(slide, gx + Inches(0.4), gy + Inches(0.65),
                  gw - Inches(0.8), Inches(0.015), fill=LIGHT_GRAY)
        _tb(slide, gx + Inches(0.2), gy + Inches(0.8), gw - Inches(0.4), Inches(1.6),
            desc, sz=11, color=SUB_TEXT, align=PP_ALIGN.CENTER)

    _page_num(slide, 4, TOTAL)

# ============================================================
# Slide 5: 核心功能
# ============================================================
def s05_features():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "03", "核心功能模块", "P0 核心 + P1 增强")

    modules = [
        ("AI 智能对话", DANGER, "P0", [
            "自然语言交互", "Function Calling 意图路由",
            "个性化记忆系统", "多轮上下文理解"]),
        ("课表 & 成绩", TENC_BLUE, "P0", [
            "周视图课程表", "成绩查询 & GPA",
            "学分进度追踪", "真实教务系统对接"]),
        ("DDL 管理", SUCCESS, "P0", [
            "作业/考试截止日期", "优先级排序",
            "逾期/紧急提醒", "日历视图"]),
        ("食堂导航", WARN, "P0", [
            "6 大食堂实时状态", "窗口导览 & 分类",
            "AI 推荐今天吃什么", "校区筛选"]),
        ("校车时刻", PURPLE, "P0", [
            "望江/江安/华西线路", "工作日/周末切换",
            "下一班倒计时", "实时状态"]),
        ("更多功能", MID_GRAY, "P1", [
            "校历查询", "选课推荐",
            "RAG 文档问答", "通知 & 天气"]),
    ]

    cw = Inches(1.9)
    ch = Inches(3.3)
    cg = Inches(0.12)
    sy = Inches(1.25)
    for i, (title, color, pri, items) in enumerate(modules):
        cx = Inches(0.4) + (cw + cg) * i
        _color_top_card(slide, cx, sy, cw, ch, color, title, items)
        # P0/P1 标签
        tag = _add_rounded(slide, cx + cw - Inches(0.55), sy + Inches(0.5),
                           Inches(0.45), Inches(0.22), fill=color)
        _tb(slide, cx + cw - Inches(0.55), sy + Inches(0.5),
            Inches(0.45), Inches(0.22), pri, sz=8, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER)

    # 用户旅程图
    _tb(slide, Inches(0.6), Inches(4.75), Inches(3), Inches(0.35),
        "典型用户旅程", sz=14, color=TENC_BLUE, bold=True)

    journey = ["打开App", "语音/文字提问", "AI识别意图", "路由到功能", "返回结果", "记忆偏好"]
    jw = Inches(1.7)
    jg = Inches(0.3)
    jy = Inches(5.15)
    for i, step in enumerate(journey):
        jx = Inches(0.6) + (jw + jg) * i
        _add_rounded(slide, jx, jy, jw, Inches(0.55), fill=TENC_BG, border=TENC_LIGHT)
        _tb(slide, jx, jy + Inches(0.08), jw, Inches(0.4),
            f"{'①②③④⑤⑥'[i]}  {step}", sz=10, color=TENC_BLUE, bold=True,
            align=PP_ALIGN.CENTER)
        if i < 5:
            # 箭头
            _tb(slide, jx + jw + Inches(0.05), jy + Inches(0.08),
                Inches(0.2), Inches(0.35), "→", sz=14, color=TENC_LIGHT,
                align=PP_ALIGN.CENTER, font="Segoe UI")

    # 底部说明
    _add_rounded(slide, Inches(0.6), Inches(5.95), Inches(11.8), Inches(0.5),
                 fill=TENC_BG, border=TENC_LIGHT)
    _tb(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.4),
        "P0 = 核心必做 (Sprint 1-4, 已基本完成)     |     P1 = 增强功能 (Sprint 5-6)     |     P2 = 未来规划",
        sz=10, color=TENC_BLUE, align=PP_ALIGN.CENTER)

    _page_num(slide, 5, TOTAL)

# ============================================================
# Slide 6: 系统架构
# ============================================================
def s06_architecture():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "04", "系统架构设计", "前后端分离 + 模块化服务 + 容器化部署")

    layers = [
        ("前端展示层", "Next.js 14 + React 19 + shadcn/ui + TailwindCSS",
         TENC_BLUE, ["登录页", "Dashboard", "课表", "成绩", "食堂", "校车", "AI 对话", "设置"]),
        ("后端服务层", "FastAPI Gateway + JWT + Redis + Alembic",
         SUCCESS, ["Auth 认证", "Academic", "Food 食堂", "Campus", "AI Intent", "Rate Limit"]),
        ("数据存储层", "持久化 + 缓存 + 外部 API",
         GOLD, ["PostgreSQL 16", "Redis 7", "教务系统 API", "通义千问 LLM"]),
    ]

    ly = Inches(1.3)
    lh = Inches(1.5)
    lg = Inches(0.35)

    for i, (name, desc, color, mods) in enumerate(layers):
        y = ly + (lh + lg) * i
        # 层背景
        _add_rounded(slide, Inches(0.5), y, Inches(12.3), lh, fill=WHITE, border=color)
        # 左侧色条
        _add_rect(slide, Inches(0.52), y + Inches(0.02), Inches(0.06), lh - Inches(0.04),
                  fill=color)
        # 层名
        _tb(slide, Inches(0.75), y + Inches(0.1), Inches(2.5), Inches(0.35),
            name, sz=14, color=color, bold=True)
        _tb(slide, Inches(0.75), y + Inches(0.45), Inches(3.2), Inches(0.25),
            desc, sz=9, color=MID_GRAY)
        # 模块方块
        mw = Inches(1.2)
        mg = Inches(0.12)
        mx_start = Inches(4.2)
        for j, mod in enumerate(mods):
            mx = mx_start + (mw + mg) * j
            _add_rounded(slide, mx, y + Inches(0.2), mw, Inches(0.55),
                         fill=TENC_BG, border=color)
            _tb(slide, mx + Inches(0.05), y + Inches(0.3), mw - Inches(0.1), Inches(0.35),
                mod, sz=9, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

        # 层间箭头
        if i < 2:
            arrow_y = y + lh + Inches(0.05)
            _tb(slide, Inches(6.2), arrow_y, Inches(0.8), Inches(0.25),
                "▼", sz=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # 右下：Docker Compose 部署图
    _tb(slide, Inches(0.6), Inches(6.1), Inches(4), Inches(0.35),
        "部署架构: Docker Compose 5 容器",
        sz=12, color=TENC_BLUE, bold=True)

    containers = [
        ("frontend", ":3000", TENC_BLUE),
        ("gateway", ":8000", SUCCESS),
        ("ai-svc", ":8005", PURPLE),
        ("postgres", ":5432", GOLD),
        ("redis", ":6379", WARN),
    ]
    dw = Inches(2.15)
    dg = Inches(0.12)
    dy = Inches(6.5)
    for i, (name, port, color) in enumerate(containers):
        dx = Inches(0.6) + (dw + dg) * i
        _add_rounded(slide, dx, dy, dw, Inches(0.55), fill=color)
        _tb(slide, dx, dy + Inches(0.08), dw, Inches(0.35),
            f"{name} {port}", sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _page_num(slide, 6, TOTAL)

# ============================================================
# Slide 7: 技术选型
# ============================================================
def s07_tech_stack():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "05", "技术选型", "生产级成熟度 | 社区活跃 | 异步优先")

    categories = [
        ("前端技术", TENC_BLUE, [
            ("Next.js 14", "SSR + App Router"),
            ("React 19", "UI 组件框架"),
            ("TailwindCSS 4", "原子化 CSS"),
            ("shadcn/ui", "高质量组件库"),
            ("Zustand", "轻量状态管理"),
            ("TanStack Query", "请求缓存"),
            ("Zod", "类型安全验证"),
        ]),
        ("后端技术", SUCCESS, [
            ("FastAPI", "异步 Python 框架"),
            ("SQLAlchemy 2.0", "异步 ORM"),
            ("Alembic", "数据库迁移"),
            ("Redis 7", "缓存/限流/会话"),
            ("PyJWT", "Token 认证"),
            ("Pydantic V2", "数据验证"),
            ("httpx", "异步 HTTP 客户端"),
        ]),
        ("AI / 数据", PURPLE, [
            ("通义千问 Qwen", "主力 LLM"),
            ("LangChain", "Function Calling"),
            ("PostgreSQL 16", "主数据库"),
            ("JSONB", "灵活数据存储"),
            ("RAG (P1)", "向量检索问答"),
            ("", ""),
            ("", ""),
        ]),
        ("DevOps", WARN, [
            ("Docker", "容器化"),
            ("Compose", "编排部署"),
            ("GitHub Actions", "CI/CD 流水线"),
            ("Ruff", "Python Lint"),
            ("ESLint", "JS/TS Lint"),
            ("pytest", "后端测试"),
            ("Vitest", "前端测试"),
        ]),
    ]

    cw = Inches(2.9)
    cg = Inches(0.2)
    cy = Inches(1.2)

    for i, (cat, color, techs) in enumerate(categories):
        cx = Inches(0.5) + (cw + cg) * i
        _white_card(slide, cx, cy, cw, Inches(5.5))
        _add_rect(slide, cx + Inches(0.02), cy + Inches(0.02),
                  cw - Inches(0.04), Inches(0.42), fill=color)
        _tb(slide, cx + Inches(0.15), cy + Inches(0.07), cw - Inches(0.3), Inches(0.35),
            cat, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        for j, (tech, desc) in enumerate(techs):
            if not tech:
                continue
            ty = cy + Inches(0.55) + Inches(0.65) * j
            _tb(slide, cx + Inches(0.15), ty, cw - Inches(0.3), Inches(0.3),
                tech, sz=12, color=DARK_TEXT, bold=True)
            _tb(slide, cx + Inches(0.15), ty + Inches(0.28), cw - Inches(0.3), Inches(0.25),
                desc, sz=9, color=SUB_TEXT)
            if j < len(techs) - 1 and techs[j+1][0]:
                _add_rect(slide, cx + Inches(0.3), ty + Inches(0.57),
                          cw - Inches(0.6), Inches(0.01), fill=LIGHT_GRAY)

    _page_num(slide, 7, TOTAL)

# ============================================================
# Slide 8: 团队分工
# ============================================================
def s08_team():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "06", "团队成员与分工", "8人团队 · 4个专业小组 · 敏捷协作")

    teams = [
        ("前端组", "2-3人", TENC_BLUE, "FE",
         ["Next.js 页面 & 路由开发", "UI/UX 组件库搭建",
          "响应式 & 移动端适配", "Zustand 状态管理", "API 对接 & 联调"]),
        ("后端组", "2-3人", SUCCESS, "BE",
         ["FastAPI 接口设计开发", "数据库建模 & 迁移",
          "教务系统 CAS 爬虫", "JWT 认证 & 限流", "单元/集成测试"]),
        ("AI 组", "1-2人", PURPLE, "AI",
         ["LLM 意图路由设计", "Function Calling 编排",
          "用户记忆系统开发", "RAG 文档问答引擎", "Prompt Engineering"]),
        ("DevOps/PM", "1人", WARN, "PM",
         ["Docker Compose 部署", "GitHub Actions CI/CD",
          "项目管理 & 进度", "Code Review 质量", "文档 & 演示准备"]),
    ]

    tw = Inches(2.85)
    tg = Inches(0.25)
    ty_start = Inches(1.25)
    th = Inches(4.0)

    for i, (name, count, color, icon, tasks) in enumerate(teams):
        tx = Inches(0.6) + (tw + tg) * i
        _white_card(slide, tx, ty_start, tw, th)
        # 头像圆
        _add_circle(slide, tx + Inches(1.0), ty_start + Inches(0.25), Inches(0.75), color)
        _tb(slide, tx + Inches(1.0), ty_start + Inches(0.35), Inches(0.75), Inches(0.55),
            icon, sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Segoe UI")
        # 名字
        _tb(slide, tx + Inches(0.1), ty_start + Inches(1.1), tw - Inches(0.2), Inches(0.3),
            name, sz=15, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, tx + Inches(0.1), ty_start + Inches(1.4), tw - Inches(0.2), Inches(0.25),
            count, sz=11, color=color, align=PP_ALIGN.CENTER)
        # 任务列表
        for j, task in enumerate(tasks):
            task_y = ty_start + Inches(1.8) + Inches(0.38) * j
            _add_circle(slide, tx + Inches(0.2), task_y + Inches(0.08),
                        Inches(0.12), color)
            _tb(slide, tx + Inches(0.4), task_y, tw - Inches(0.55), Inches(0.3),
                task, sz=10, color=SUB_TEXT)

    # 底部协作模式
    _add_rounded(slide, Inches(0.6), Inches(5.55), Inches(11.8), Inches(0.55),
                 fill=TENC_BG, border=TENC_LIGHT)
    _tb(slide, Inches(0.8), Inches(5.62), Inches(11.4), Inches(0.4),
        "协作模式:  GitHub Flow 分支管理  |  每周 Standup  |  双周 Sprint Review  |  OpenAPI 文档先行  |  Conventional Commits",
        sz=10, color=TENC_BLUE, align=PP_ALIGN.CENTER)

    # 饼图: 工作量分配
    _tb(slide, Inches(0.6), Inches(6.3), Inches(3), Inches(0.3),
        "工作量分配", sz=12, color=TENC_BLUE, bold=True)

    chart_data = ChartData()
    chart_data.categories = ['前端 35%', '后端 30%', 'AI 20%', 'DevOps 15%']
    chart_data.add_series('', (35, 30, 20, 15))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(0.4), Inches(6.3), Inches(3.5), Inches(1.15),
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)

    plot = chart.plots[0]
    colors = [TENC_BLUE, SUCCESS, PURPLE, WARN]
    for i, color in enumerate(colors):
        point = plot.series[0].points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color

    _page_num(slide, 8, TOTAL)

# ============================================================
# Slide 9: 开发计划
# ============================================================
def s09_timeline():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "07", "开发计划与里程碑", "16 周敏捷迭代")

    # 甘特图
    sprints = [
        ("S0 基础设施", 0, 2, MID_GRAY, "Git/Docker/CI"),
        ("S1 核心骨架", 2, 4, TENC_BLUE, "登录/布局/DB"),
        ("S2 学业MVP", 4, 6, SUCCESS, "课表+成绩+教务"),
        ("S3 食堂+校车", 6, 8, WARN, "食堂导航+校车"),
        ("S4 AI 对话", 8, 10, PURPLE, "意图路由+FC"),
        ("S5 P1 扩展", 10, 12, GOLD, "选课+RAG"),
        ("S6 打磨优化", 12, 14, CYAN, "UI优化+深色"),
        ("S7 测试交付", 14, 16, DANGER, "测试/文档/Demo"),
    ]

    chart_left = Inches(2.8)
    chart_w = Inches(10.0)
    label_x = Inches(0.5)
    row_h = Inches(0.55)
    sy = Inches(1.4)

    # 周标签
    for w in range(17):
        wx = chart_left + int(chart_w * w / 16)
        if w < 16:
            _tb(slide, wx, sy - Inches(0.25), Inches(0.6), Inches(0.2),
                f"W{w+1}", sz=7, color=MID_GRAY, align=PP_ALIGN.CENTER, font="Segoe UI")
        # 竖虚线
        if w % 2 == 0:
            _add_rect(slide, wx + Inches(0.02), sy, Emu(6000), row_h * 8,
                      fill=RGBColor(0xE8,0xE8,0xEE))

    for i, (name, start, end, color, desc) in enumerate(sprints):
        y = sy + row_h * i
        # 标签
        _tb(slide, label_x, y + Inches(0.05), Inches(2.2), Inches(0.35),
            name, sz=10, color=DARK_TEXT, bold=True)
        # 条
        bx = chart_left + int(chart_w * start / 16)
        bw = int(chart_w * (end - start) / 16)
        _add_rounded(slide, bx, y + Inches(0.08), bw, Inches(0.35), fill=color)
        _tb(slide, bx + Inches(0.05), y + Inches(0.1), bw - Inches(0.1), Inches(0.3),
            desc, sz=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 当前进度线
    cur_w = 6
    cur_x = chart_left + int(chart_w * cur_w / 16)
    _add_rect(slide, cur_x, sy - Inches(0.1), Emu(18000), row_h * 8 + Inches(0.3), fill=DANGER)
    _add_rounded(slide, cur_x - Inches(0.25), sy - Inches(0.35),
                 Inches(0.6), Inches(0.22), fill=DANGER)
    _tb(slide, cur_x - Inches(0.25), sy - Inches(0.35),
        Inches(0.6), Inches(0.22), "当前", sz=8, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)

    # 里程碑
    ms_y = sy + row_h * 8 + Inches(0.25)
    milestones = [
        (4, "骨架验收", GOLD), (8, "MVP完成", SUCCESS),
        (10, "AI上线", PURPLE), (16, "最终交付", DANGER),
    ]
    for week, label, color in milestones:
        mx = chart_left + int(chart_w * week / 16) - Inches(0.15)
        # 菱形
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, mx, ms_y, Inches(0.3), Inches(0.3))
        d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()
        _tb(slide, mx - Inches(0.35), ms_y + Inches(0.32), Inches(1.0), Inches(0.25),
            f"W{week} {label}", sz=8, color=color, bold=True, align=PP_ALIGN.CENTER)

    # 进度统计
    _tb(slide, Inches(0.6), Inches(6.3), Inches(3), Inches(0.3),
        "Sprint 进度分布", sz=12, color=TENC_BLUE, bold=True)

    chart_data = CategoryChartData()
    chart_data.categories = ['S0','S1','S2','S3','S4','S5','S6','S7']
    chart_data.add_series('计划工时(人天)', (8, 16, 20, 16, 20, 16, 12, 12))
    chart_data.add_series('已完成(人天)', (8, 16, 20, 16, 4, 0, 0, 0))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(6.3), Inches(6.0), Inches(1.15),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(8)
    chart.legend.include_in_layout = False
    s0 = chart.series[0]; s0.format.fill.solid(); s0.format.fill.fore_color.rgb = TENC_BG
    s1 = chart.series[1]; s1.format.fill.solid(); s1.format.fill.fore_color.rgb = TENC_BLUE

    _page_num(slide, 9, TOTAL)

# ============================================================
# Slide 10: 技术难点与风险
# ============================================================
def s10_risks():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "08", "技术难点与风险应对")

    risks = [
        ("教务系统对接", "高", DANGER,
         "• 反爬机制 + 验证码\n• Session 管理复杂\n• 接口可能变更",
         "• 双MD5加密还原\n• Mock模式并行开发\n• 定期监控接口"),
        ("AI 意图准确率", "中", WARN,
         "• 自然语言歧义\n• 多意图混合\n• FC 稳定性",
         "• Prompt Engineering\n• 意图分类优化\n• 页面导航降级兜底"),
        ("性能与并发", "中", WARN,
         "• 高峰期并发\n• 数据库慢查询\n• LLM 响应延迟",
         "• Redis 多级缓存\n• 数据库索引优化\n• SSE 流式输出"),
        ("团队协作", "低", SUCCESS,
         "• 接口对齐\n• 风格不一致\n• 进度延期",
         "• OpenAPI 文档先行\n• CI 自动 Lint\n• 敏捷迭代+周同步"),
    ]

    rw = Inches(2.85)
    rg = Inches(0.25)
    ry = Inches(1.2)
    rh = Inches(4.2)

    for i, (title, level, color, challenge, solution) in enumerate(risks):
        rx = Inches(0.6) + (rw + rg) * i
        _white_card(slide, rx, ry, rw, rh)
        # 标题 + 标签
        _tb(slide, rx + Inches(0.15), ry + Inches(0.12), Inches(1.6), Inches(0.3),
            title, sz=13, color=DARK_TEXT, bold=True)
        _add_rounded(slide, rx + rw - Inches(0.8), ry + Inches(0.12),
                     Inches(0.65), Inches(0.25), fill=color)
        _tb(slide, rx + rw - Inches(0.8), ry + Inches(0.12),
            Inches(0.65), Inches(0.25), f"风险{level}", sz=9, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER)

        # 挑战
        _tb(slide, rx + Inches(0.15), ry + Inches(0.55), rw - Inches(0.3), Inches(0.22),
            "挑战", sz=10, color=DANGER, bold=True)
        _tb(slide, rx + Inches(0.15), ry + Inches(0.8), rw - Inches(0.3), Inches(1.3),
            challenge, sz=9, color=SUB_TEXT)

        # 分割线
        _add_rect(slide, rx + Inches(0.3), ry + Inches(2.15),
                  rw - Inches(0.6), Inches(0.015), fill=LIGHT_GRAY)

        # 应对
        _tb(slide, rx + Inches(0.15), ry + Inches(2.3), rw - Inches(0.3), Inches(0.22),
            "应对方案", sz=10, color=SUCCESS, bold=True)
        _tb(slide, rx + Inches(0.15), ry + Inches(2.55), rw - Inches(0.3), Inches(1.3),
            solution, sz=9, color=SUB_TEXT)

    # 风险矩阵图
    _tb(slide, Inches(0.6), Inches(5.6), Inches(3), Inches(0.3),
        "风险矩阵", sz=12, color=TENC_BLUE, bold=True)

    # 用图表做风险矩阵
    chart_data = ChartData()
    chart_data.categories = ['教务对接', 'AI准确率', '性能并发', '团队协作']
    chart_data.add_series('影响度', (9, 7, 6, 4))
    chart_data.add_series('概率', (7, 5, 5, 3))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4), Inches(5.6), Inches(5.5), Inches(1.5),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(8)
    s0 = chart.series[0]; s0.format.fill.solid(); s0.format.fill.fore_color.rgb = DANGER
    s1 = chart.series[1]; s1.format.fill.solid(); s1.format.fill.fore_color.rgb = WARN

    # 缓存策略表
    _tb(slide, Inches(6.5), Inches(5.6), Inches(3), Inches(0.3),
        "缓存策略", sz=12, color=TENC_BLUE, bold=True)
    cache_items = [
        ("课表/校历/校车", "Redis 24h", SUCCESS),
        ("食堂营业状态", "实时计算", TENC_BLUE),
        ("LLM 回复", "Redis 1h (Hash)", PURPLE),
        ("用户会话", "Redis 30min/7d", WARN),
    ]
    for i, (data, strategy, color) in enumerate(cache_items):
        cy = Inches(5.95) + Inches(0.28) * i
        _add_circle(slide, Inches(6.5), cy + Inches(0.06), Inches(0.12), color)
        _tb(slide, Inches(6.75), cy, Inches(2.2), Inches(0.25),
            data, sz=9, color=DARK_TEXT, bold=True)
        _tb(slide, Inches(9.0), cy, Inches(2), Inches(0.25),
            strategy, sz=9, color=SUB_TEXT)

    _page_num(slide, 10, TOTAL)

# ============================================================
# Slide 11: 当前进展
# ============================================================
def s11_progress():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "09", "项目当前进展", "P0 核心模块完成度 60%")

    # 总进度条
    _tb(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.3),
        "整体完成度", sz=14, color=TENC_BLUE, bold=True)
    _progress_bar(slide, Inches(2.8), Inches(1.25), Inches(9.8), Inches(0.35), 0.6, SUCCESS)
    _tb(slide, Inches(2.8), Inches(1.27), Inches(5.88), Inches(0.3),
        "60%  P0 核心模块已完成", sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 模块进度
    _tb(slide, Inches(0.6), Inches(1.85), Inches(4), Inches(0.3),
        "各模块完成度", sz=13, color=TENC_BLUE, bold=True)

    module_progress = [
        ("认证系统", 1.0, SUCCESS),
        ("Dashboard", 1.0, SUCCESS),
        ("课程表", 1.0, SUCCESS),
        ("DDL 管理", 0.6, WARN),
        ("食堂导航", 0.9, SUCCESS),
        ("校车时刻", 0.9, SUCCESS),
        ("AI 对话", 0.2, DANGER),
        ("设置页面", 0.8, TENC_BLUE),
    ]

    for i, (name, ratio, color) in enumerate(module_progress):
        my = Inches(2.2) + Inches(0.38) * i
        _tb(slide, Inches(0.8), my, Inches(1.3), Inches(0.3),
            name, sz=10, color=DARK_TEXT, bold=True)
        _progress_bar(slide, Inches(2.2), my + Inches(0.05), Inches(3.0), Inches(0.22),
                      ratio, color)
        _tb(slide, Inches(5.35), my, Inches(0.6), Inches(0.3),
            f"{int(ratio*100)}%", sz=9, color=color, bold=True, font="Segoe UI")

    # 已完成列表
    _white_card(slide, Inches(6.2), Inches(1.8), Inches(3.3), Inches(3.8))
    _add_rect(slide, Inches(6.22), Inches(1.82), Inches(3.26), Inches(0.35), fill=SUCCESS)
    _tb(slide, Inches(6.4), Inches(1.85), Inches(3), Inches(0.3),
        "✓  已完成", sz=12, color=WHITE, bold=True)
    done_items = [
        "认证系统 (真实教务登录+JWT)",
        "Dashboard 仪表盘",
        "课程表 (周视图+颜色区分)",
        "DDL 管理 (本地CRUD)",
        "食堂导航 (6食堂+窗口)",
        "校车时刻 (三校区+倒计时)",
        "AI 对话 UI 骨架",
        "CI/CD + Docker 部署",
    ]
    for i, item in enumerate(done_items):
        _tb(slide, Inches(6.4), Inches(2.25) + Inches(0.35) * i,
            Inches(3), Inches(0.3), f"✓ {item}", sz=9, color=SUB_TEXT)

    # 待完成列表
    _white_card(slide, Inches(9.7), Inches(1.8), Inches(3.2), Inches(3.8))
    _add_rect(slide, Inches(9.72), Inches(1.82), Inches(3.16), Inches(0.35), fill=WARN)
    _tb(slide, Inches(9.9), Inches(1.85), Inches(3), Inches(0.3),
        "○  待完成", sz=12, color=WHITE, bold=True)
    todo_items = [
        "DDL 后端 API 对接",
        "AI 接入通义千问 LLM",
        "深色模式主题切换",
        "成绩详情独立页面",
        "校历页面",
        "通知系统",
        "数据库完整迁移",
        "用户记忆 & 推荐",
    ]
    for i, item in enumerate(todo_items):
        _tb(slide, Inches(9.9), Inches(2.25) + Inches(0.35) * i,
            Inches(3), Inches(0.3), f"○ {item}", sz=9, color=SUB_TEXT)

    # 代码统计
    _tb(slide, Inches(0.6), Inches(5.55), Inches(3), Inches(0.3),
        "代码统计", sz=13, color=TENC_BLUE, bold=True)

    code_stats = [
        ("前端代码", "~3,500 行", TENC_BLUE),
        ("后端代码", "~2,800 行", SUCCESS),
        ("配置/CI", "~500 行", WARN),
        ("总计", "~6,800 行", DARK_TEXT),
    ]
    for i, (label, lines, color) in enumerate(code_stats):
        _stat_card(slide, Inches(0.6) + Inches(2.4) * i, Inches(5.9),
                   Inches(2.2), Inches(0.9), lines, label, color)

    # 提交统计图
    _tb(slide, Inches(6.2), Inches(5.55), Inches(3), Inches(0.3),
        "Git 提交趋势", sz=13, color=TENC_BLUE, bold=True)

    chart_data = CategoryChartData()
    chart_data.categories = ['W1-2', 'W3-4', 'W5-6']
    chart_data.add_series('提交数', (12, 28, 35))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(6.0), Inches(5.8), Inches(3.5), Inches(1.2),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    s0 = chart.series[0]; s0.format.fill.solid(); s0.format.fill.fore_color.rgb = TENC_BLUE

    # 测试覆盖率
    _tb(slide, Inches(9.8), Inches(5.55), Inches(3), Inches(0.3),
        "测试覆盖率", sz=13, color=TENC_BLUE, bold=True)

    chart_data2 = ChartData()
    chart_data2.categories = ['后端', '前端']
    chart_data2.add_series('已覆盖', (65, 40))
    chart_data2.add_series('未覆盖', (35, 60))

    chart_shape2 = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        Inches(9.6), Inches(5.8), Inches(3.3), Inches(1.2),
        chart_data2
    )
    chart2 = chart_shape2.chart
    chart2.has_legend = True
    chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2.legend.font.size = Pt(7)
    s2a = chart2.series[0]; s2a.format.fill.solid(); s2a.format.fill.fore_color.rgb = SUCCESS
    s2b = chart2.series[1]; s2b.format.fill.solid(); s2b.format.fill.fore_color.rgb = LIGHT_GRAY

    _page_num(slide, 11, TOTAL)

# ============================================================
# Slide 12: 总结亮点
# ============================================================
def s12_highlights():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "10", "项目亮点与展望")

    # 左列：亮点
    _tb(slide, Inches(0.8), Inches(1.2), Inches(4), Inches(0.35),
        "⭐  项目亮点", sz=16, color=GOLD, bold=True)

    highlights = [
        ("真实数据对接", "直连川大教务系统，双MD5+CAS认证", TENC_BLUE),
        ("AI 原生设计", "LLM Function Calling 驱动核心交互", PURPLE),
        ("现代化全栈", "Next.js 14 + FastAPI + Docker 容器化", SUCCESS),
        ("体验优先", "深色主题、响应式、移动端适配", WARN),
        ("工程化实践", "CI/CD、Code Review、自动化测试", DANGER),
    ]
    for i, (title, desc, color) in enumerate(highlights):
        hy = Inches(1.7) + Inches(0.85) * i
        _white_card(slide, Inches(0.6), hy, Inches(5.8), Inches(0.72))
        _add_rect(slide, Inches(0.62), hy + Inches(0.02), Inches(0.06), Inches(0.68), fill=color)
        _tb(slide, Inches(0.9), hy + Inches(0.08), Inches(2.5), Inches(0.28),
            title, sz=13, color=DARK_TEXT, bold=True)
        _tb(slide, Inches(0.9), hy + Inches(0.38), Inches(5.2), Inches(0.25),
            desc, sz=10, color=SUB_TEXT)

    # 右列：下一步
    _tb(slide, Inches(6.9), Inches(1.2), Inches(5), Inches(0.35),
        "🚀  下一步计划", sz=16, color=TENC_BLUE, bold=True)

    next_steps = [
        ("1", "AI 对话全面接入", "通义千问 + Function Calling 上线"),
        ("2", "DDL 系统完善", "后端 API + 数据库 + 提醒通知"),
        ("3", "个性化记忆系统", "口味偏好、课程习惯、智能推荐"),
        ("4", "RAG 文档问答", "校规、选课手册知识库检索"),
        ("5", "全面测试与上线", "性能优化、安全加固、用户测试"),
    ]
    for i, (num, title, desc) in enumerate(next_steps):
        ny = Inches(1.7) + Inches(0.85) * i
        _white_card(slide, Inches(6.9), ny, Inches(5.8), Inches(0.72), border=TENC_LIGHT)
        _add_circle(slide, Inches(7.05), ny + Inches(0.15), Inches(0.4), TENC_BLUE)
        _tb(slide, Inches(7.05), ny + Inches(0.18), Inches(0.4), Inches(0.35),
            num, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Segoe UI")
        _tb(slide, Inches(7.6), ny + Inches(0.08), Inches(4.8), Inches(0.28),
            title, sz=13, color=DARK_TEXT, bold=True)
        _tb(slide, Inches(7.6), ny + Inches(0.38), Inches(4.8), Inches(0.25),
            desc, sz=10, color=SUB_TEXT)

    # 底部预期成果
    _add_rounded(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0),
                 fill=TENC_BG, border=TENC_BLUE)
    _tb(slide, Inches(0.9), Inches(6.1), Inches(3), Inches(0.3),
        "预期交付成果", sz=13, color=TENC_BLUE, bold=True)
    deliverables = [
        "完整可运行的 Web 应用 (含前后端+AI)",
        "Docker 一键部署方案",
        "完整技术文档 + API 文档",
        "用户测试报告 + 演示视频",
    ]
    for i, d in enumerate(deliverables):
        dx = Inches(0.9) + Inches(3.0) * i
        _tb(slide, dx, Inches(6.42), Inches(2.8), Inches(0.4),
            f"✓ {d}", sz=9, color=SUB_TEXT)

    _page_num(slide, 12, TOTAL)

# ============================================================
# Slide 13: 致谢
# ============================================================
def s13_thanks():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 封面同款背景
    _add_rect(slide, 0, 0, SW, SH, fill=TENC_DARK)
    _add_rect(slide, 0, 0, Inches(6.5), SH, fill=TENC_BLUE)

    c1 = _add_circle(slide, Inches(10), Inches(-0.5), Inches(4), RGBColor(0x1A, 0x3D, 0x7A))

    _tb(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.0),
        "谢谢！", sz=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
        "SCU Assistant — 让校园生活更智能",
        sz=22, color=RGBColor(0xBB,0xD5,0xFF), align=PP_ALIGN.CENTER)

    _add_rect(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.03), fill=TENC_LIGHT)

    _tb(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.4),
        "四川大学计算机学院  |  2026 春季  |  软件工程课程设计",
        sz=13, color=MID_GRAY, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
        "欢迎提问与交流", sz=14, color=RGBColor(0x88,0xAA,0xDD),
        align=PP_ALIGN.CENTER)

    _page_num(slide, 13, TOTAL)

# ============================================================
# Slide 14: 附录 — 数据库 ER 概览
# ============================================================
def s14_appendix():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=NEAR_WHITE)
    _section_title(slide, "附", "附录：数据库表设计", "核心数据模型")

    tables = [
        ("users", "用户表", ["student_id", "name", "campus", "major", "grade", "preferences (JSONB)"], TENC_BLUE),
        ("schedules", "课程表", ["user_id", "course_name", "teacher", "location", "weekday", "weeks"], SUCCESS),
        ("deadlines", "DDL", ["user_id", "title", "due_date", "priority", "status", "source"], WARN),
        ("canteens", "食堂", ["name", "campus", "open_time", "close_time", "meal_type"], GOLD),
        ("bus_schedules", "校车", ["route", "departure", "arrival", "time", "is_weekend"], PURPLE),
        ("conversations", "对话", ["user_id", "title", "created_at"], DANGER),
        ("messages", "消息", ["conversation_id", "role", "content", "tool_calls (JSONB)"], CYAN),
        ("user_memories", "记忆", ["user_id", "key", "value", "category", "confidence"], SCU_RED),
    ]

    tw = Inches(2.95)
    tg = Inches(0.2)
    rows_per_col = 4

    for i, (tname, tdesc, fields, color) in enumerate(tables):
        col = i // rows_per_col
        row = i % rows_per_col
        tx = Inches(0.5) + col * (Inches(6.5) + Inches(0.3))
        ty = Inches(1.2) + row * Inches(1.45)

        _white_card(slide, tx, ty, Inches(6.3), Inches(1.3))
        _add_rect(slide, tx + Inches(0.02), ty + Inches(0.02),
                  Inches(0.06), Inches(1.26), fill=color)
        _tb(slide, tx + Inches(0.2), ty + Inches(0.08), Inches(2), Inches(0.25),
            tname, sz=12, color=color, bold=True, font="Consolas")
        _tb(slide, tx + Inches(2.0), ty + Inches(0.08), Inches(1.5), Inches(0.25),
            tdesc, sz=10, color=SUB_TEXT)
        # 字段
        field_text = "  |  ".join(fields)
        _tb(slide, tx + Inches(0.2), ty + Inches(0.4), Inches(5.8), Inches(0.8),
            field_text, sz=8, color=MID_GRAY, font="Consolas")

    _page_num(slide, 14, TOTAL)

# ============================================================
# 生成
# ============================================================
print("生成 PPT 中...")
s01_cover()
s02_toc()
s03_background()
s04_positioning()
s05_features()
s06_architecture()
s07_tech_stack()
s08_team()
s09_timeline()
s10_risks()
s11_progress()
s12_highlights()
s13_thanks()
s14_appendix()

out = os.path.join(os.path.dirname(__file__), "SCU_Assistant_开题报告_v3.pptx")
prs.save(out)
print(f"✓ 已生成: {out}")
print(f"  共 {len(prs.slides)} 页")
