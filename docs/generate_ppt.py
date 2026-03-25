"""
SCU Assistant 项目开题报告 PPT 生成脚本
风格参考：互联网大厂（字节/腾讯/阿里）项目汇报模板
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# ============================================================
# 品牌色系 — 川大红 + 科技蓝
# ============================================================
SCU_RED = RGBColor(0xC4, 0x12, 0x30)
SCU_RED_LIGHT = RGBColor(0xE8, 0x17, 0x3A)
SCU_GOLD = RGBColor(0xD4, 0xA8, 0x43)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG2 = RGBColor(0x16, 0x21, 0x3E)
TECH_BLUE = RGBColor(0x0F, 0x3D, 0x6B)
ACCENT_BLUE = RGBColor(0x1E, 0x90, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
NEAR_BLACK = RGBColor(0x22, 0x22, 0x33)
GREEN = RGBColor(0x00, 0xC9, 0xA7)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
PURPLE = RGBColor(0xA8, 0x5C, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ============================================================
# 工具函数
# ============================================================
def add_bg(slide, color=DARK_BG):
    """深色背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_accent_bar(slide, left, top, width, height, color=SCU_RED):
    """强调色装饰条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=16, color=LIGHT_GRAY, bold=False,
             space_before=Pt(4), space_after=Pt(4), alignment=PP_ALIGN.LEFT,
             font_name="Microsoft YaHei"):
    """在已有文本框中追加段落"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_card(slide, left, top, width, height, fill_color=DARK_BG2):
    """卡片背景"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
    card.line.width = Pt(1)
    return card


def add_icon_circle(slide, left, top, size, color, text, font_size=20):
    """圆形图标"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Segoe UI Emoji"
    return circle


def slide_number_footer(slide, num, total):
    """页码"""
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12


# ============================================================
# Slide 1: 封面
# ============================================================
def make_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BG)

    # 顶部装饰线
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    # 左侧竖条
    add_accent_bar(slide, Inches(1.2), Inches(2.0), Inches(0.08), Inches(2.8), SCU_GOLD)

    # 主标题
    add_text_box(slide, Inches(1.6), Inches(2.0), Inches(10), Inches(1.2),
                 "SCU Assistant", font_size=52, color=WHITE, bold=True,
                 font_name="Segoe UI")
    add_text_box(slide, Inches(1.6), Inches(3.0), Inches(10), Inches(0.8),
                 "四川大学智能校园助手", font_size=36, color=SCU_GOLD, bold=True)
    add_text_box(slide, Inches(1.6), Inches(3.9), Inches(10), Inches(0.6),
                 "项目开题报告  |  软件工程课程设计", font_size=20, color=LIGHT_GRAY)

    # 底部信息
    add_accent_bar(slide, 0, Inches(6.4), SLIDE_W, Inches(0.02), RGBColor(0x33, 0x33, 0x55))
    add_text_box(slide, Inches(1.6), Inches(6.6), Inches(5), Inches(0.5),
                 "四川大学计算机学院  |  2026年春季学期", font_size=14, color=MID_GRAY)
    add_text_box(slide, Inches(7), Inches(6.6), Inches(5), Inches(0.5),
                 "汇报日期：2026年3月", font_size=14, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)

    # 右侧装饰圆
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(1.0),
                                 Inches(2.5), Inches(2.5))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(0xC4, 0x12, 0x30)
    c1.line.fill.background()
    # 内圆
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(1.3),
                                 Inches(1.9), Inches(1.9))
    c2.fill.solid()
    c2.fill.fore_color.rgb = DARK_BG
    c2.line.fill.background()
    add_text_box(slide, Inches(10.8), Inches(1.7), Inches(1.9), Inches(1.2),
                 "SCU\nAI+", font_size=28, color=SCU_GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    slide_number_footer(slide, 1, TOTAL_SLIDES)


# ============================================================
# Slide 2: 目录
# ============================================================
def make_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(1.2), Inches(0.6), Inches(4), Inches(0.7),
                 "CONTENTS", font_size=14, color=MID_GRAY, bold=True,
                 font_name="Segoe UI")
    add_text_box(slide, Inches(1.2), Inches(1.0), Inches(4), Inches(0.7),
                 "目录", font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.2), Inches(1.7), Inches(0.8), Inches(0.05), SCU_RED)

    toc_items = [
        ("01", "项目背景与痛点", "校园信息碎片化现状"),
        ("02", "产品定位与目标", "AI驱动的一站式校园助手"),
        ("03", "核心功能模块", "课表/DDL/食堂/校车/AI对话"),
        ("04", "系统架构设计", "前后端分离 + 微服务"),
        ("05", "技术选型", "Next.js + FastAPI + LLM"),
        ("06", "团队分工", "前端/后端/AI/DevOps"),
        ("07", "开发计划 & 里程碑", "16周敏捷迭代"),
        ("08", "技术难点与风险", "教务对接/AI准确率"),
        ("09", "项目现状与演示", "已完成核心P0功能"),
        ("10", "总结与展望", "下一步计划"),
    ]

    start_y = Inches(2.2)
    col1_x = Inches(1.5)
    col2_x = Inches(7.0)

    for i, (num, title, desc) in enumerate(toc_items):
        col_x = col1_x if i < 5 else col2_x
        row_y = start_y + Inches(0.9) * (i % 5)

        # 序号
        add_text_box(slide, col_x, row_y, Inches(0.6), Inches(0.5),
                     num, font_size=28, color=SCU_RED, bold=True,
                     font_name="Segoe UI")
        # 标题
        add_text_box(slide, col_x + Inches(0.7), row_y, Inches(4), Inches(0.35),
                     title, font_size=18, color=WHITE, bold=True)
        # 描述
        add_text_box(slide, col_x + Inches(0.7), row_y + Inches(0.35), Inches(4), Inches(0.3),
                     desc, font_size=12, color=MID_GRAY)

    slide_number_footer(slide, 2, TOTAL_SLIDES)


# ============================================================
# Slide 3: 项目背景与痛点
# ============================================================
def make_background(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "01", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "项目背景与痛点分析", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    # 痛点卡片
    pain_points = [
        ("信息碎片化", "课表查教务、校车看公众号\n食堂信息靠口口相传\n学生需在5+个平台间切换",
         ORANGE, "!"),
        ("系统体验差", "教务系统 UI 落后\n移动端适配差\n操作路径长、加载慢",
         SCU_RED_LIGHT, "X"),
        ("缺乏智能化", "无个性化推荐\n无 DDL 提醒\n无统一问答入口",
         PURPLE, "?"),
        ("数据不互通", "各系统数据孤岛\n无法交叉查询\n信息时效性差",
         ACCENT_BLUE, "#"),
    ]

    card_w = Inches(2.6)
    card_h = Inches(3.8)
    start_x = Inches(0.8)
    card_y = Inches(1.8)
    gap = Inches(0.3)

    for i, (title, desc, color, icon) in enumerate(pain_points):
        x = start_x + (card_w + gap) * i
        add_card(slide, x, card_y, card_w, card_h)

        # 图标
        add_icon_circle(slide, x + Inches(0.9), card_y + Inches(0.3),
                        Inches(0.7), color, icon, font_size=24)

        add_text_box(slide, x + Inches(0.2), card_y + Inches(1.2), Inches(2.2), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), card_y + Inches(1.7), Inches(2.2), Inches(2.0),
                     desc, font_size=13, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # 右侧结论
    add_card(slide, Inches(11.6), card_y, Inches(1.5), card_h, SCU_RED)
    txBox = add_text_box(slide, Inches(11.65), card_y + Inches(0.5),
                         Inches(1.4), Inches(3.0),
                         "", font_size=16, color=WHITE, bold=True,
                         alignment=PP_ALIGN.CENTER)
    tf = txBox.text_frame
    tf.paragraphs[0].text = "亟需"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_para(tf, "一站式", font_size=20, color=SCU_GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_para(tf, "智能", font_size=20, color=SCU_GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_para(tf, "校园", font_size=20, color=SCU_GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_para(tf, "平台", font_size=20, color=SCU_GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)

    # 底部数据条
    add_card(slide, Inches(0.8), Inches(5.9), Inches(11.8), Inches(1.0))
    stats = [
        ("50,000+", "川大在校学生"),
        ("5+", "日常使用平台数"),
        ("70%", "学生反馈查课不便"),
        ("0", "现有智能助手产品"),
    ]
    for i, (num, label) in enumerate(stats):
        sx = Inches(1.5) + Inches(2.9) * i
        add_text_box(slide, sx, Inches(6.0), Inches(2), Inches(0.45),
                     num, font_size=26, color=SCU_GOLD, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Segoe UI")
        add_text_box(slide, sx, Inches(6.4), Inches(2), Inches(0.35),
                     label, font_size=12, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

    slide_number_footer(slide, 3, TOTAL_SLIDES)


# ============================================================
# Slide 4: 产品定位与目标
# ============================================================
def make_positioning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "02", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "产品定位与核心目标", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    # 核心定位 大卡片
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.2), TECH_BLUE)
    add_text_box(slide, Inches(1.2), Inches(1.8), Inches(4.5), Inches(0.5),
                 "产品定位", font_size=14, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(4.8), Inches(0.5),
                 "AI 驱动的一站式智能校园助手", font_size=22, color=WHITE, bold=True)
    txBox = add_text_box(slide, Inches(1.2), Inches(2.8), Inches(4.8), Inches(0.8),
                         "", font_size=14, color=LIGHT_GRAY)
    tf = txBox.text_frame
    tf.paragraphs[0].text = "自然语言对话 + 传统页面导航 双模式交互"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    add_para(tf, "对接真实教务数据，提供个性化智能服务", font_size=14, color=LIGHT_GRAY)

    # 目标用户
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.2))
    add_text_box(slide, Inches(7.2), Inches(1.8), Inches(4), Inches(0.5),
                 "目标用户", font_size=14, color=ACCENT_CYAN, bold=True)
    users = [
        "四川大学全日制本科生 & 研究生",
        "日均使用教务/食堂/校车等校园服务的学生",
        "期望通过 AI 提升信息获取效率的用户",
    ]
    y_off = Inches(2.3)
    for u in users:
        add_text_box(slide, Inches(7.5), y_off, Inches(5), Inches(0.4),
                     f"   {u}", font_size=14, color=LIGHT_GRAY)
        # bullet
        add_accent_bar(slide, Inches(7.2), y_off + Inches(0.1), Inches(0.12), Inches(0.12), GREEN)
        y_off += Inches(0.4)

    # 四个目标卡片
    goals = [
        ("效率提升", "一句话查课表、查成绩\n减少 80% 操作步骤", GREEN, ">"),
        ("智能推荐", "AI 意图识别 + 个性化\n食堂/选课智能推荐", ACCENT_BLUE, "*"),
        ("数据整合", "教务/食堂/校车/DDL\n一站式信息聚合", SCU_GOLD, "+"),
        ("体验升级", "现代化 UI/UX 设计\n移动端优先、深色主题", PURPLE, "~"),
    ]

    card_w = Inches(2.7)
    gap = Inches(0.3)
    start_x = Inches(0.8)
    card_y = Inches(4.2)

    for i, (title, desc, color, icon) in enumerate(goals):
        x = start_x + (card_w + gap) * i
        add_card(slide, x, card_y, card_w, Inches(2.8))
        add_icon_circle(slide, x + Inches(0.15), card_y + Inches(0.2),
                        Inches(0.55), color, icon, font_size=20)
        add_text_box(slide, x + Inches(0.8), card_y + Inches(0.25), Inches(1.8), Inches(0.4),
                     title, font_size=17, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), card_y + Inches(1.0), Inches(2.3), Inches(1.5),
                     desc, font_size=13, color=LIGHT_GRAY)

    slide_number_footer(slide, 4, TOTAL_SLIDES)


# ============================================================
# Slide 5: 核心功能模块
# ============================================================
def make_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "03", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "核心功能模块", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    features = [
        ("AI 智能对话", "自然语言交互\nFunction Calling 意图路由\n个性化记忆系统\n多轮上下文理解",
         SCU_RED_LIGHT, "P0"),
        ("课表 & 成绩", "周视图课程表\n成绩查询 & GPA\n学分进度追踪\n对接真实教务系统",
         ACCENT_BLUE, "P0"),
        ("DDL 管理", "作业/考试截止日期\n优先级排序\n逾期/紧急提醒\n日历视图",
         GREEN, "P0"),
        ("食堂导航", "6大食堂实时状态\n窗口导览 & 分类\nAI 推荐今天吃什么\n校区筛选",
         ORANGE, "P0"),
        ("校车时刻", "三校区线路\n工作日/周末时刻表\n下一班倒计时\n实时状态",
         PURPLE, "P0"),
        ("更多功能", "校历查询\n选课推荐 (P1)\nRAG 文档问答 (P1)\n通知系统 (P1)",
         MID_GRAY, "P1"),
    ]

    card_w = Inches(1.85)
    card_h = Inches(4.8)
    start_x = Inches(0.6)
    card_y = Inches(1.6)
    gap = Inches(0.18)

    for i, (title, desc, color, priority) in enumerate(features):
        x = start_x + (card_w + gap) * i
        add_card(slide, x, card_y, card_w, card_h)

        # 顶部色条
        add_accent_bar(slide, x, card_y, card_w, Inches(0.06), color)

        # 优先级标签
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(1.15), card_y + Inches(0.2),
                                      Inches(0.55), Inches(0.28))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tag_tf = tag.text_frame
        tag_tf.paragraphs[0].text = priority
        tag_tf.paragraphs[0].font.size = Pt(10)
        tag_tf.paragraphs[0].font.color.rgb = WHITE
        tag_tf.paragraphs[0].font.bold = True
        tag_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 标题
        add_text_box(slide, x + Inches(0.1), card_y + Inches(0.6), Inches(1.65), Inches(0.5),
                     title, font_size=15, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # 分割线
        add_accent_bar(slide, x + Inches(0.3), card_y + Inches(1.15),
                       Inches(1.25), Inches(0.02), RGBColor(0x33, 0x33, 0x55))

        # 描述文字
        lines = desc.split("\n")
        for j, line in enumerate(lines):
            add_text_box(slide, x + Inches(0.12), card_y + Inches(1.4) + Inches(0.4) * j,
                         Inches(1.6), Inches(0.35),
                         f"  {line}", font_size=11, color=LIGHT_GRAY)
            # 小圆点
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          x + Inches(0.15),
                                          card_y + Inches(1.5) + Inches(0.4) * j,
                                          Inches(0.08), Inches(0.08))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

    # 底部说明
    add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
                 "P0 = 核心必做 (Sprint 1-4)    |    P1 = 增强功能 (Sprint 5-6)    |    P2 = 未来规划",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    slide_number_footer(slide, 5, TOTAL_SLIDES)


# ============================================================
# Slide 6: 系统架构
# ============================================================
def make_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "04", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "系统架构设计", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    # 三层架构
    layers = [
        ("前端展示层", "Next.js 14 + React 19 + shadcn/ui + TailwindCSS",
         [("登录页", ACCENT_BLUE), ("Dashboard", GREEN), ("课表/成绩", ORANGE),
          ("食堂", SCU_GOLD), ("校车", PURPLE), ("AI 对话", SCU_RED_LIGHT), ("设置", MID_GRAY)],
         Inches(1.5), ACCENT_BLUE),

        ("后端服务层", "FastAPI + SQLAlchemy + JWT + Redis",
         [("Auth 认证", GREEN), ("Academic", ACCENT_BLUE), ("Food 食堂", ORANGE),
          ("Campus 校车", PURPLE), ("AI Intent", SCU_RED_LIGHT), ("Rate Limit", MID_GRAY)],
         Inches(3.4), GREEN),

        ("数据存储层", "PostgreSQL + Redis + 外部API",
         [("PostgreSQL", ACCENT_BLUE), ("Redis Cache", ORANGE), ("教务系统", SCU_RED_LIGHT),
          ("LLM API", PURPLE)],
         Inches(5.3), SCU_GOLD),
    ]

    for layer_name, layer_desc, modules, y_pos, layer_color in layers:
        # 层背景
        add_card(slide, Inches(0.5), y_pos, Inches(12.3), Inches(1.6))

        # 层标题
        add_accent_bar(slide, Inches(0.5), y_pos, Inches(0.06), Inches(1.6), layer_color)
        add_text_box(slide, Inches(0.8), y_pos + Inches(0.1), Inches(2.5), Inches(0.35),
                     layer_name, font_size=16, color=layer_color, bold=True)
        add_text_box(slide, Inches(0.8), y_pos + Inches(0.45), Inches(3), Inches(0.3),
                     layer_desc, font_size=10, color=MID_GRAY)

        # 模块方块
        mod_start_x = Inches(3.8)
        mod_w = Inches(1.25)
        mod_gap = Inches(0.15)
        for j, (mod_name, mod_color) in enumerate(modules):
            mx = mod_start_x + (mod_w + mod_gap) * j
            mod_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              mx, y_pos + Inches(0.3),
                                              mod_w, Inches(0.9))
            mod_box.fill.solid()
            mod_box.fill.fore_color.rgb = RGBColor(0x22, 0x2A, 0x45)
            mod_box.line.color.rgb = mod_color
            mod_box.line.width = Pt(1.5)
            add_text_box(slide, mx + Inches(0.05), y_pos + Inches(0.5),
                         Inches(1.15), Inches(0.4),
                         mod_name, font_size=10, color=WHITE,
                         alignment=PP_ALIGN.CENTER)

    # 层间箭头
    for arrow_y in [Inches(3.15), Inches(5.05)]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                        Inches(6.2), arrow_y, Inches(0.4), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

    slide_number_footer(slide, 6, TOTAL_SLIDES)


# ============================================================
# Slide 7: 技术选型
# ============================================================
def make_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "05", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "技术选型", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    categories = [
        ("前端技术", ACCENT_BLUE, [
            ("Next.js 14", "React 框架，SSR + App Router"),
            ("TailwindCSS 4", "原子化 CSS 工具"),
            ("shadcn/ui", "高质量无头组件库"),
            ("Zustand", "轻量状态管理"),
            ("TanStack Query", "数据请求 & 缓存"),
        ]),
        ("后端技术", GREEN, [
            ("FastAPI", "异步 Python Web 框架"),
            ("SQLAlchemy 2.0", "异步 ORM + Alembic 迁移"),
            ("Redis 7", "缓存 / 限流 / 会话"),
            ("JWT (PyJWT)", "认证令牌管理"),
            ("Pydantic V2", "数据验证 & 序列化"),
        ]),
        ("AI / 数据", PURPLE, [
            ("通义千问 Qwen", "主力 LLM 模型"),
            ("LangChain", "Function Calling 编排"),
            ("PostgreSQL 16", "主数据库 (JSONB)"),
            ("RAG (P1)", "向量检索文档问答"),
        ]),
        ("DevOps", ORANGE, [
            ("Docker Compose", "容器化部署编排"),
            ("GitHub Actions", "CI/CD 自动化"),
            ("Ruff + ESLint", "代码质量检查"),
            ("pytest + Vitest", "前后端测试"),
        ]),
    ]

    card_w = Inches(2.85)
    card_h = Inches(5.0)
    start_x = Inches(0.6)
    card_y = Inches(1.5)
    gap = Inches(0.2)

    for i, (cat_name, cat_color, techs) in enumerate(categories):
        x = start_x + (card_w + gap) * i
        add_card(slide, x, card_y, card_w, card_h)
        add_accent_bar(slide, x, card_y, card_w, Inches(0.05), cat_color)

        add_text_box(slide, x + Inches(0.2), card_y + Inches(0.2), Inches(2.5), Inches(0.4),
                     cat_name, font_size=17, color=cat_color, bold=True)

        for j, (tech_name, tech_desc) in enumerate(techs):
            ty = card_y + Inches(0.75) + Inches(0.85) * j
            # 技术名
            add_text_box(slide, x + Inches(0.2), ty, Inches(2.5), Inches(0.3),
                         tech_name, font_size=14, color=WHITE, bold=True)
            # 描述
            add_text_box(slide, x + Inches(0.2), ty + Inches(0.3), Inches(2.5), Inches(0.3),
                         tech_desc, font_size=11, color=MID_GRAY)

    slide_number_footer(slide, 7, TOTAL_SLIDES)


# ============================================================
# Slide 8: 团队分工
# ============================================================
def make_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "06", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "团队成员与分工", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    # 团队概览
    add_card(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8))
    add_text_box(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(0.55),
                 "8人团队  |  4个专业小组  |  敏捷协作  |  每周 Standup + 双周 Sprint Review",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    teams = [
        ("前端组", "2-3人", ACCENT_BLUE, "FE",
         ["Next.js 页面开发", "UI/UX 组件实现", "响应式适配", "状态管理 & API 对接"]),
        ("后端组", "2-3人", GREEN, "BE",
         ["FastAPI 接口开发", "数据库设计 & 迁移", "教务系统爬虫对接", "认证 & 限流中间件"]),
        ("AI 组", "1-2人", PURPLE, "AI",
         ["LLM 意图路由", "Function Calling", "用户记忆系统", "RAG 文档问答"]),
        ("DevOps/PM", "1人", ORANGE, "PM",
         ["Docker 部署", "CI/CD 流水线", "项目管理 & 文档", "代码 Review & 质量"]),
    ]

    card_w = Inches(2.7)
    card_h = Inches(4.2)
    start_x = Inches(0.8)
    card_y = Inches(2.6)
    gap = Inches(0.3)

    for i, (name, count, color, icon, tasks) in enumerate(teams):
        x = start_x + (card_w + gap) * i
        add_card(slide, x, card_y, card_w, card_h)

        # 头像圆
        add_icon_circle(slide, x + Inches(0.95), card_y + Inches(0.3),
                        Inches(0.75), color, icon, font_size=22)

        add_text_box(slide, x + Inches(0.1), card_y + Inches(1.2), Inches(2.5), Inches(0.35),
                     name, font_size=18, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), card_y + Inches(1.55), Inches(2.5), Inches(0.3),
                     count, font_size=13, color=color,
                     alignment=PP_ALIGN.CENTER)

        # 职责列表
        for j, task in enumerate(tasks):
            ty = card_y + Inches(2.0) + Inches(0.42) * j
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), ty + Inches(0.08),
                                          Inches(0.08), Inches(0.08))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            add_text_box(slide, x + Inches(0.4), ty, Inches(2.2), Inches(0.35),
                         task, font_size=12, color=LIGHT_GRAY)

    slide_number_footer(slide, 8, TOTAL_SLIDES)


# ============================================================
# Slide 9: 开发计划 (甘特图风格)
# ============================================================
def make_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "07", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "开发计划 & 里程碑", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    # 甘特图区域
    chart_left = Inches(2.8)
    chart_right = Inches(12.5)
    chart_width = chart_right - chart_left
    label_x = Inches(0.5)
    row_h = Inches(0.55)
    start_y = Inches(1.8)

    # 周标签 (W1-W16)
    for w in range(16):
        wx = chart_left + int(chart_width * w / 16)
        add_text_box(slide, wx, start_y - Inches(0.3), Inches(0.6), Inches(0.25),
                     f"W{w+1}", font_size=8, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER, font_name="Segoe UI")
        # 竖线
        if w % 2 == 0:
            vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            wx + Inches(0.25), start_y,
                                            Emu(8000), Inches(4.8))
            vline.fill.solid()
            vline.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
            vline.line.fill.background()

    sprints = [
        ("Sprint 0: 基础设施", 0, 2, RGBColor(0x4A, 0x5A, 0x7A), "Git/Docker/CI"),
        ("Sprint 1: 核心骨架", 2, 4, ACCENT_BLUE, "登录/布局/DB"),
        ("Sprint 2: 学业 MVP", 4, 6, GREEN, "课表+成绩+教务对接"),
        ("Sprint 3: 食堂+校车", 6, 8, ORANGE, "食堂导航+校车时刻"),
        ("Sprint 4: AI 对话", 8, 10, PURPLE, "意图路由+Function Call"),
        ("Sprint 5: P1 扩展", 10, 12, SCU_GOLD, "选课推荐+RAG"),
        ("Sprint 6: 打磨优化", 12, 14, ACCENT_CYAN, "晨间简报+UI优化"),
        ("Sprint 7: 测试交付", 14, 16, SCU_RED_LIGHT, "Bug修复+文档+演示"),
    ]

    for i, (name, start_w, end_w, color, milestone) in enumerate(sprints):
        y = start_y + row_h * i

        # 标签
        add_text_box(slide, label_x, y + Inches(0.05), Inches(2.3), Inches(0.35),
                     name, font_size=11, color=WHITE, bold=True)

        # 进度条
        bar_x = chart_left + int(chart_width * start_w / 16)
        bar_w = int(chart_width * (end_w - start_w) / 16)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      bar_x, y + Inches(0.08), bar_w, Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 里程碑文字
        add_text_box(slide, bar_x + Inches(0.1), y + Inches(0.1),
                     bar_w - Inches(0.1), Inches(0.3),
                     milestone, font_size=9, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

    # 里程碑标记
    milestones_mark = [
        (4, "骨架验收", SCU_GOLD),
        (8, "MVP 完成", GREEN),
        (10, "AI 上线", PURPLE),
        (16, "最终交付", SCU_RED_LIGHT),
    ]
    ms_y = start_y + row_h * 8 + Inches(0.3)
    for week, label, color in milestones_mark:
        mx = chart_left + int(chart_width * week / 16) - Inches(0.1)
        # 菱形标记
        diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, mx, ms_y, Inches(0.25), Inches(0.25))
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = color
        diamond.line.fill.background()
        add_text_box(slide, mx - Inches(0.4), ms_y + Inches(0.3), Inches(1.2), Inches(0.3),
                     f"W{week} {label}", font_size=10, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # 当前进度指示
    current_week = 6  # 假设当前第6周
    cur_x = chart_left + int(chart_width * current_week / 16)
    cur_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       cur_x, start_y - Inches(0.1),
                                       Emu(18000), Inches(5.3))
    cur_line.fill.solid()
    cur_line.fill.fore_color.rgb = SCU_RED
    cur_line.line.fill.background()
    add_text_box(slide, cur_x - Inches(0.3), start_y - Inches(0.45), Inches(0.8), Inches(0.25),
                 "当前", font_size=9, color=SCU_RED, bold=True,
                 alignment=PP_ALIGN.CENTER)

    slide_number_footer(slide, 9, TOTAL_SLIDES)


# ============================================================
# Slide 10: 技术难点与风险 + 饼图
# ============================================================
def make_risks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "08", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "技术难点与风险应对", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    risks = [
        ("教务系统对接", "高",
         "反爬机制、验证码识别\nSession 管理、接口变更",
         "双 MD5 加密还原\nMock 模式开发测试\n定期监控接口变化",
         SCU_RED_LIGHT),
        ("AI 意图准确率", "中",
         "自然语言歧义\n多意图混合\nFunction Call 稳定性",
         "Prompt Engineering\n意图分类预训练\n兜底页面导航降级",
         PURPLE),
        ("性能与并发", "中",
         "高峰期并发访问\n数据库慢查询\nLLM 响应延迟",
         "Redis 多级缓存\n数据库索引优化\nSSE 流式输出",
         ORANGE),
        ("团队协作", "低",
         "前后端接口对齐\n代码风格不一致\n进度延期风险",
         "OpenAPI 文档先行\nCI 自动化检查\n敏捷迭代 + 每周同步",
         ACCENT_BLUE),
    ]

    card_w = Inches(2.85)
    card_h = Inches(4.5)
    start_x = Inches(0.6)
    card_y = Inches(1.6)
    gap = Inches(0.2)

    risk_colors = {"高": SCU_RED_LIGHT, "中": ORANGE, "低": GREEN}

    for i, (title, level, challenge, solution, color) in enumerate(risks):
        x = start_x + (card_w + gap) * i
        add_card(slide, x, card_y, card_w, card_h)

        # 风险等级
        lvl_color = risk_colors[level]
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(2.05), card_y + Inches(0.15),
                                      Inches(0.6), Inches(0.28))
        tag.fill.solid()
        tag.fill.fore_color.rgb = lvl_color
        tag.line.fill.background()
        tag_tf = tag.text_frame
        tag_tf.paragraphs[0].text = f"风险{level}"
        tag_tf.paragraphs[0].font.size = Pt(9)
        tag_tf.paragraphs[0].font.color.rgb = WHITE
        tag_tf.paragraphs[0].font.bold = True
        tag_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, x + Inches(0.15), card_y + Inches(0.15), Inches(2), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True)

        # 挑战
        add_text_box(slide, x + Inches(0.15), card_y + Inches(0.65), Inches(2.5), Inches(0.25),
                     "挑战", font_size=11, color=color, bold=True)
        add_text_box(slide, x + Inches(0.15), card_y + Inches(0.9), Inches(2.5), Inches(1.3),
                     challenge, font_size=11, color=LIGHT_GRAY)

        # 分割线
        add_accent_bar(slide, x + Inches(0.3), card_y + Inches(2.3),
                       Inches(2.2), Inches(0.02), RGBColor(0x33, 0x33, 0x55))

        # 应对方案
        add_text_box(slide, x + Inches(0.15), card_y + Inches(2.5), Inches(2.5), Inches(0.25),
                     "应对方案", font_size=11, color=GREEN, bold=True)
        add_text_box(slide, x + Inches(0.15), card_y + Inches(2.8), Inches(2.5), Inches(1.5),
                     solution, font_size=11, color=LIGHT_GRAY)

    # 底部工作量分配饼图
    add_card(slide, Inches(0.6), Inches(6.3), Inches(11.8), Inches(0.9))
    add_text_box(slide, Inches(1.0), Inches(6.4), Inches(2), Inches(0.3),
                 "工作量分配", font_size=13, color=WHITE, bold=True)
    # 水平条形比例
    segments = [
        ("前端 35%", 0.35, ACCENT_BLUE),
        ("后端 30%", 0.30, GREEN),
        ("AI 20%", 0.20, PURPLE),
        ("DevOps 15%", 0.15, ORANGE),
    ]
    bar_start = Inches(3.5)
    bar_total_w = Inches(8.5)
    bar_y = Inches(6.45)
    bar_h = Inches(0.5)
    cumulative = 0
    for label, ratio, color in segments:
        seg_x = bar_start + int(bar_total_w * cumulative)
        seg_w = int(bar_total_w * ratio)
        seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, seg_x, bar_y, seg_w, bar_h)
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.line.fill.background()
        add_text_box(slide, seg_x + Inches(0.05), bar_y + Inches(0.08),
                     seg_w - Inches(0.1), Inches(0.35),
                     label, font_size=11, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        cumulative += ratio

    slide_number_footer(slide, 10, TOTAL_SLIDES)


# ============================================================
# Slide 11: 项目现状与进度
# ============================================================
def make_status(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "09", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "项目当前进展", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    # 完成进度总览
    add_card(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2))
    add_text_box(slide, Inches(1.2), Inches(1.6), Inches(3), Inches(0.4),
                 "整体完成度", font_size=16, color=WHITE, bold=True)
    # 进度条背景
    prog_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(4.5), Inches(1.7), Inches(7.5), Inches(0.45))
    prog_bg.fill.solid()
    prog_bg.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
    prog_bg.line.fill.background()
    # 进度条
    prog = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4.5), Inches(1.7), Inches(4.5), Inches(0.45))
    prog.fill.solid()
    prog.fill.fore_color.rgb = GREEN
    prog.line.fill.background()
    add_text_box(slide, Inches(4.5), Inches(1.72), Inches(4.5), Inches(0.4),
                 "60%  P0 核心模块已完成", font_size=14, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 已完成 & 进行中
    completed = [
        "认证系统 (教务真实登录 + JWT)",
        "Dashboard 仪表盘 (今日课程/成绩/学分)",
        "课程表页面 (周视图 + 颜色区分)",
        "DDL 管理页面 (本地 CRUD + 优先级)",
        "食堂导航 (6食堂 + 窗口导览)",
        "校车时刻 (三校区 + 倒计时)",
        "AI 对话 UI 骨架",
        "CI/CD 流水线",
    ]

    in_progress = [
        "DDL 后端 API 对接",
        "AI 接入通义千问 LLM",
        "深色模式主题切换",
        "成绩详情独立页面",
        "校历页面",
        "通知系统",
        "数据库完整迁移",
    ]

    # 已完成列
    add_card(slide, Inches(0.8), Inches(3.0), Inches(5.7), Inches(4.2))
    add_text_box(slide, Inches(1.1), Inches(3.1), Inches(3), Inches(0.4),
                 "已完成", font_size=16, color=GREEN, bold=True)
    for i, item in enumerate(completed):
        iy = Inches(3.55) + Inches(0.42) * i
        # 勾选框
        check = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(1.1), iy, Inches(0.22), Inches(0.22))
        check.fill.solid()
        check.fill.fore_color.rgb = GREEN
        check.line.fill.background()
        add_text_box(slide, Inches(1.0), iy - Inches(0.04), Inches(0.42), Inches(0.3),
                     " \u2713", font_size=12, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.5), iy, Inches(4.8), Inches(0.3),
                     item, font_size=13, color=LIGHT_GRAY)

    # 进行中列
    add_card(slide, Inches(6.8), Inches(3.0), Inches(5.7), Inches(4.2))
    add_text_box(slide, Inches(7.1), Inches(3.1), Inches(3), Inches(0.4),
                 "进行中 / 待完成", font_size=16, color=ORANGE, bold=True)
    for i, item in enumerate(in_progress):
        iy = Inches(3.55) + Inches(0.42) * i
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.1), iy + Inches(0.05),
                                      Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE
        dot.line.fill.background()
        add_text_box(slide, Inches(7.45), iy, Inches(4.8), Inches(0.3),
                     item, font_size=13, color=LIGHT_GRAY)

    slide_number_footer(slide, 11, TOTAL_SLIDES)


# ============================================================
# Slide 12: 总结与展望
# ============================================================
def make_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), SCU_RED)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
                 "10", font_size=40, color=SCU_RED, bold=True, font_name="Segoe UI")
    add_text_box(slide, Inches(1.8), Inches(0.55), Inches(6), Inches(0.5),
                 "总结与展望", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), SCU_RED)

    # 项目亮点
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(3.0))
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(4), Inches(0.4),
                 "项目亮点", font_size=18, color=SCU_GOLD, bold=True)

    highlights = [
        ("真实数据对接", "非 Mock 数据，直连川大教务系统"),
        ("AI 原生设计", "LLM Function Calling 驱动核心交互"),
        ("现代化技术栈", "Next.js 14 + FastAPI + Docker 全栈"),
        ("用户体验优先", "深色主题、响应式、移动端适配"),
        ("工程化实践", "CI/CD、代码审查、自动化测试"),
    ]
    for i, (title, desc) in enumerate(highlights):
        hy = Inches(2.2) + Inches(0.45) * i
        add_accent_bar(slide, Inches(1.2), hy + Inches(0.08),
                       Inches(0.12), Inches(0.12), SCU_GOLD)
        add_text_box(slide, Inches(1.5), hy, Inches(2), Inches(0.3),
                     title, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, Inches(3.3), hy, Inches(3), Inches(0.3),
                     desc, font_size=12, color=LIGHT_GRAY)

    # 下一步计划
    add_card(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(3.0))
    add_text_box(slide, Inches(7.4), Inches(1.7), Inches(4), Inches(0.4),
                 "下一步计划", font_size=18, color=ACCENT_CYAN, bold=True)

    next_steps = [
        "完成 AI 对话模块接入通义千问",
        "DDL 系统后端 API 全面对接",
        "用户记忆系统 & 个性化推荐",
        "RAG 文档问答引擎",
        "全面测试 & 性能优化",
    ]
    for i, step in enumerate(next_steps):
        sy = Inches(2.2) + Inches(0.45) * i
        # 序号圆
        add_icon_circle(slide, Inches(7.3), sy, Inches(0.3), ACCENT_CYAN,
                        str(i + 1), font_size=12)
        add_text_box(slide, Inches(7.75), sy, Inches(4.5), Inches(0.3),
                     step, font_size=13, color=LIGHT_GRAY)

    # 底部致谢
    add_card(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0), TECH_BLUE)
    add_text_box(slide, Inches(1.2), Inches(5.2), Inches(11), Inches(0.7),
                 "THANK YOU", font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Segoe UI")
    add_text_box(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.5),
                 "SCU Assistant  —  让校园生活更智能", font_size=18, color=SCU_GOLD,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.2), Inches(6.4), Inches(11), Inches(0.4),
                 "四川大学计算机学院  |  2026 春季  |  软件工程课程设计",
                 font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    slide_number_footer(slide, 12, TOTAL_SLIDES)


# ============================================================
# 生成 PPT
# ============================================================
print("正在生成 PPT...")

make_cover(prs)
make_toc(prs)
make_background(prs)
make_positioning(prs)
make_features(prs)
make_architecture(prs)
make_tech_stack(prs)
make_team(prs)
make_timeline(prs)
make_risks(prs)
make_status(prs)
make_summary(prs)

output_path = os.path.join(os.path.dirname(__file__), "SCU_Assistant_开题报告.pptx")
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"共 {len(prs.slides)} 页")
