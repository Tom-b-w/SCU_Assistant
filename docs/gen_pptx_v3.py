"""
SCU Assistant 开题报告 PPT v3 — 基于答辩风.pptx模板
直接复用模板的背景图片、校徽、配色和布局风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(SCRIPT_DIR, "_template_imgs")

# ============================================================
# 模板图片路径
# ============================================================
BG_COVER   = os.path.join(IMG_DIR, "s01_img0.jpg")   # 封面/结尾全屏背景 13.3x7.5"
LOGO_BIG   = os.path.join(IMG_DIR, "s01_img1.png")   # 左上角校徽 2.4x0.7"
BG_CONTENT = os.path.join(IMG_DIR, "s02_img2.jpg")   # 内容页背景 11.2x7.5"
BG_SECTION = os.path.join(IMG_DIR, "s03_img3.jpg")   # section分隔页左图 4.2x7.5"
LOGO_SMALL = os.path.join(IMG_DIR, "s04_img4.png")   # 右上角小logo 2.4x0.7"

# ============================================================
# 配色 — 与模板一致
# ============================================================
SCU_RED      = RGBColor(0xC4, 0x12, 0x30)
SCU_RED_DARK = RGBColor(0x8B, 0x0E, 0x22)
GOLD         = RGBColor(0xD4, 0xA8, 0x43)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE   = RGBColor(0xF5, 0xF5, 0xF8)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY     = RGBColor(0x88, 0x88, 0x99)
DARK_TEXT     = RGBColor(0x33, 0x33, 0x33)
SUB_TEXT     = RGBColor(0x66, 0x66, 0x77)
ACCENT_BLUE  = RGBColor(0x1E, 0x90, 0xFF)
GREEN        = RGBColor(0x00, 0xC9, 0xA7)
ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
TEAL         = RGBColor(0x00, 0x96, 0x88)

# ---- 创建演示文稿（使用模板的slide尺寸）----
prs = Presentation()
prs.slide_width  = 12192000   # 与模板一致
prs.slide_height = 6858000
SW = prs.slide_width
SH = prs.slide_height

# ============================================================
# 工具函数
# ============================================================
def _rect(slide, left, top, w, h, fill=None, border=None, border_w=Pt(0.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    return s

def _rrect(slide, left, top, w, h, fill=None, border=None, border_w=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    return s

def _circle(slide, left, top, d, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def _line(slide, left, top, w):
    """红色横线（模板标题下方装饰线）"""
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = SCU_RED
    s.line.fill.background()
    return s

def _tb(slide, left, top, w, h, text, sz=14, color=DARK_TEXT,
        bold=False, italic=False, font="Microsoft YaHei",
        align=PP_ALIGN.LEFT, line_spacing=1.15):
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    p.line_spacing = Pt(sz * line_spacing)
    tf.auto_size = None
    return tf

def _multi_tb(slide, left, top, w, h, lines, default_sz=12, default_color=DARK_TEXT,
              align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(lines):
        text = item[0]
        sz = item[1] if len(item) > 1 else default_sz
        color = item[2] if len(item) > 2 else default_color
        bold = item[3] if len(item) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.alignment = align
        p.line_spacing = Pt(sz * 1.4)
    return tf

def _add_overlay(slide, color, alpha_val):
    """添加半透明遮罩层"""
    from pptx.oxml.ns import qn
    from lxml import etree
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = color
    # 通过 XML 设置透明度
    spPr = overlay._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        # fill 可能在其他位置
        ln = spPr
        for child in spPr:
            if 'solidFill' in child.tag:
                solidFill = child
                break
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', alpha_val)
    overlay.line.fill.background()

# ---- 模板页面组件 ----

def _add_content_bg(slide):
    """内容页：模板背景图（右侧）+ 白色遮罩（淡化）+ 右上角小logo"""
    slide.shapes.add_picture(BG_CONTENT, Inches(2.08), 0,
                             Inches(11.25), Inches(7.5))
    # 白色半透明遮罩 — 淡化背景让文字清晰（模板原版也有此处理）
    _add_overlay(slide, WHITE, '82000')  # 82%不透明白色，背景隐约可见
    # 右上角小logo
    slide.shapes.add_picture(LOGO_SMALL, Inches(10.70), Inches(0.23),
                             Inches(2.4), Inches(0.7))

def _add_section_bg(slide):
    """Section分隔页：左侧装饰图"""
    # 深色背景
    _rect(slide, 0, 0, SW, SH, fill=RGBColor(0xF0, 0xE8, 0xE4))
    slide.shapes.add_picture(BG_SECTION, Inches(1.60), 0,
                             Inches(4.22), Inches(7.5))

def _title_with_line(slide, x, y, num, title):
    """模板风格标题：编号+标题 + 下方红线（与模板S4一致）"""
    _tb(slide, x, y, Inches(5), Inches(0.55),
        f"{num} {title}", sz=26, color=DARK_TEXT, bold=True)
    _line(slide, x - Inches(0.05), y + Inches(0.55), Inches(2.63))

def _subtitle(slide, x, y, text):
    """标题下方副标题"""
    _tb(slide, x, y, Inches(10), Inches(0.5),
        text, sz=16, color=SCU_RED, bold=True)


# ============================================================
# S01 — 封面（复用模板封面背景+校徽）
# ============================================================
def s01_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 模板封面背景
    slide.shapes.add_picture(BG_COVER, Inches(-0.02), 0,
                             Inches(13.35), Inches(7.5))
    # 半透明深色遮罩（让文字更清晰）
    _add_overlay(slide, RGBColor(0x00, 0x00, 0x00), '30000')  # 30%黑色遮罩
    # 左上角校徽
    slide.shapes.add_picture(LOGO_BIG, Inches(0.15), Inches(0.15),
                             Inches(2.4), Inches(0.7))
    # 主标题
    _tb(slide, Inches(2.0), Inches(2.0), Inches(9.5), Inches(1.0),
        "SCU Assistant", sz=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
        font="Segoe UI")
    _tb(slide, Inches(2.0), Inches(3.0), Inches(9.5), Inches(0.8),
        "四川大学智能校园助手", sz=40, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    # 副标题
    _tb(slide, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.6),
        "软件工程课程设计  ·  开题报告",
        sz=22, color=WHITE, align=PP_ALIGN.CENTER)
    # 团队
    _tb(slide, Inches(2.0), Inches(5.0), Inches(9.5), Inches(0.5),
        "谭博文  徐锐学  孔垂骄  张傲楚  毛立业  朱圣相  李亚飞  谭旭睿  覃泽锴",
        sz=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# S02 — 目录（复用模板TOC布局）
# ============================================================
def s02_toc():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 左侧红色区域
    _rect(slide, 0, 0, SW, SH, fill=SCU_RED)
    # 右侧背景图 + 白色遮罩淡化
    slide.shapes.add_picture(BG_CONTENT, Inches(2.08), 0,
                             Inches(11.25), Inches(7.5))
    # 在右侧区域加白色遮罩，让文字清晰
    from pptx.oxml.ns import qn
    from lxml import etree
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(2.08), 0, Inches(11.25), SH)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = WHITE
    spPr = overlay._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', '78000')  # 78%不透明白色
    overlay.line.fill.background()

    # 左侧"目录"文字
    _tb(slide, Inches(0.30), Inches(4.0), Inches(2.3), Inches(0.5),
        "CONTENTS", sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0.30), Inches(2.7), Inches(2.3), Inches(1.2),
        "目 录", sz=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 目录条目 — 字体放大
    items = [
        ("01", "问题与出发点"),
        ("02", "需求分析"),
        ("03", "项目概述与功能"),
        ("04", "团队成员与分工"),
        ("05", "技术路线"),
        ("06", "展示结果与未来计划"),
    ]
    for i, (num, title) in enumerate(items):
        y = Inches(1.10) + Inches(i * 1.06)
        _tb(slide, Inches(5.5), y, Inches(1.2), Inches(0.7),
            num, sz=28, color=SCU_RED, bold=True, font="Segoe UI")
        _tb(slide, Inches(6.8), y + Inches(0.08), Inches(5.0), Inches(0.6),
            title, sz=22, color=DARK_TEXT, bold=True)


# ============================================================
# S03 — Part.01 问题与出发点
# ============================================================
def s03_sec():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_bg(slide)
    # 右侧深色区域
    _rect(slide, Inches(5.8), 0, Inches(7.5), SH, fill=RGBColor(0x2A, 0x0A, 0x14))
    # Part 编号（模板位置: 6.96, 2.80）
    _tb(slide, Inches(6.96), Inches(2.80), Inches(1.3), Inches(0.5),
        "Part.01", sz=18, color=GOLD, bold=True, font="Segoe UI")
    # 标题（模板位置: 5.68, 3.31）
    _tb(slide, Inches(5.68), Inches(3.31), Inches(6.1), Inches(1.0),
        "问题与出发点", sz=36, color=WHITE, bold=True)
    # 右边窄条
    _rect(slide, Inches(13.08), 0, Inches(0.25), SH, fill=SCU_RED)


# ============================================================
# S04 — 1.1 选题背景
# ============================================================
def s04_background():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 白色底 + 内容页背景
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "1.1", "选题背景")
    _subtitle(slide, Inches(1.16), Inches(1.20), "校园信息化现状与痛点")

    # 三列卡片（模板S4有三个卡片区域）
    cards = [
        ("01", "信息碎片化",
         "课表、成绩、校车、食堂分布在 5+ 个独立平台，学生需在多个系统间切换",
         SCU_RED),
        ("02", "系统体验差",
         "教务系统 UI 陈旧，移动端适配差，查课表需要 5+ 步操作",
         ORANGE),
        ("03", "缺乏智能化",
         "无 AI 问答入口，无个性化推荐，无统一的信息聚合能力",
         PURPLE),
    ]

    # 三列区域（模板: 1.16/4.94/8.72, y=2.56, w=3.6, h=3.8）
    for i, (num, title, desc, color) in enumerate(cards):
        x = Inches(1.16) + Inches(i * 3.78)
        y = Inches(2.20)
        # 卡片背景
        _rrect(slide, x, y, Inches(3.6), Inches(4.2),
               fill=NEAR_WHITE, border=color, border_w=Pt(1.5))
        # 编号
        _circle(slide, x + Inches(1.4), y + Inches(0.3), Inches(0.7), color)
        _tb(slide, x + Inches(1.42), y + Inches(0.38), Inches(0.7), Inches(0.5),
            num, sz=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
            font="Segoe UI")
        # 标题
        _tb(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.0), Inches(0.5),
            title, sz=20, color=color, bold=True, align=PP_ALIGN.CENTER)
        # 描述
        _tb(slide, x + Inches(0.3), y + Inches(1.8), Inches(3.0), Inches(2.0),
            desc, sz=14, color=SUB_TEXT, align=PP_ALIGN.CENTER, line_spacing=1.6)

    # 底部结论条
    _rrect(slide, Inches(1.16), Inches(6.6), Inches(11.0), Inches(0.5),
           fill=SCU_RED)
    _tb(slide, Inches(1.16), Inches(6.63), Inches(11.0), Inches(0.45),
        "市场空白：四川大学尚无一站式智能校园助手，6万+师生的刚需",
        sz=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# S05 — 1.2 数据支撑
# ============================================================
def s05_data():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "1.2", "数据来源")
    _subtitle(slide, Inches(1.16), Inches(1.20), "项目核心数据获取方式")

    # 四个数据来源卡片（2x2网格）
    sources = [
        ("教务系统", SCU_RED,
         "四川大学教务处综合教学管理系统",
         "通过 CAS 统一认证登录，获取\n课表、成绩、学分、GPA 等数据"),
        ("校园公开信息", ACCENT_BLUE,
         "学校官网 / 后勤服务 / 校车中心",
         "食堂名称与窗口信息、校车发车\n时刻表（望江/江安/华西）"),
        ("用户自主录入", GREEN,
         "用户在系统内创建和管理",
         "DDL 截止日期、作业/考试提醒\n个人待办事项与优先级标记"),
        ("AI 大模型", PURPLE,
         "通义千问 API（阿里云）",
         "自然语言问答、意图识别、智能\n推荐，为对话模式提供 AI 能力"),
    ]

    for i, (title, color, source, desc) in enumerate(sources):
        row = i // 2
        col = i % 2
        x = Inches(0.91) + Inches(col * 5.7)
        y = Inches(1.90) + Inches(row * 2.35)
        card_w = Inches(5.4)
        card_h = Inches(2.1)

        _rrect(slide, x, y, card_w, card_h, fill=WHITE, border=color)
        # 左侧色条
        _rect(slide, x, y, Inches(0.08), card_h, fill=color)
        # 编号圆
        _circle(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.6), color)
        _tb(slide, x + Inches(0.27), y + Inches(0.33), Inches(0.6), Inches(0.45),
            str(i + 1), sz=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
            font="Segoe UI")
        # 标题
        _tb(slide, x + Inches(1.0), y + Inches(0.2), Inches(3.5), Inches(0.4),
            title, sz=18, color=color, bold=True)
        # 数据源
        _tb(slide, x + Inches(1.0), y + Inches(0.65), Inches(4.0), Inches(0.35),
            source, sz=13, color=SUB_TEXT, italic=True)
        # 描述
        _tb(slide, x + Inches(1.0), y + Inches(1.10), Inches(4.0), Inches(0.85),
            desc, sz=13, color=DARK_TEXT, line_spacing=1.5)


# ============================================================
# S06 — Part.02 需求分析
# ============================================================
def s06_sec():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_bg(slide)
    _rect(slide, Inches(5.8), 0, Inches(7.5), SH, fill=RGBColor(0x2A, 0x0A, 0x14))
    _tb(slide, Inches(6.96), Inches(2.80), Inches(1.3), Inches(0.5),
        "Part.02", sz=18, color=GOLD, bold=True, font="Segoe UI")
    _tb(slide, Inches(5.68), Inches(3.31), Inches(6.1), Inches(1.0),
        "需求分析", sz=36, color=WHITE, bold=True)
    _rect(slide, Inches(13.08), 0, Inches(0.25), SH, fill=SCU_RED)


# ============================================================
# S07 — 2.1 目标用户与核心需求
# ============================================================
def s07_requirement():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "2.1", "目标用户与核心需求")

    # 左侧：目标用户（模板S6样式）
    _subtitle(slide, Inches(1.16), Inches(1.20), "目标用户")

    users = [
        "四川大学全日制本科生 & 研究生（6万+）",
        "日均需查课表 / 校车 / 食堂的在校学生",
        "期望通过 AI 提升日常效率的用户群体",
    ]
    for i, u in enumerate(users):
        y = Inches(1.80) + Inches(i * 0.60)
        _circle(slide, Inches(1.16), y + Inches(0.05), Inches(0.35), SCU_RED)
        _tb(slide, Inches(1.65), y, Inches(4.5), Inches(0.45),
            u, sz=14, color=DARK_TEXT)

    # 右侧：核心需求
    _tb(slide, Inches(6.5), Inches(1.20), Inches(5.5), Inches(0.45),
        "核心功能需求", sz=18, color=SCU_RED, bold=True)

    needs = [
        ("效率提升", "一句话查课表，减少 80% 操作步骤", SCU_RED),
        ("智能交互", "AI 自然语言理解 + 意图识别 + 个性化", PURPLE),
        ("数据整合", "教务/食堂/校车/DDL 一站式聚合", ACCENT_BLUE),
        ("实时同步", "直连教务系统，数据即时更新", GREEN),
        ("体验升级", "现代 UI + 深色主题 + 移动端优先", ORANGE),
        ("任务管理", "DDL 截止提醒 + 优先级排序 + 日历视图", TEAL),
    ]
    for i, (title, desc, color) in enumerate(needs):
        y = Inches(1.80) + Inches(i * 0.75)
        _rect(slide, Inches(6.5), y, Inches(0.08), Inches(0.55), fill=color)
        _tb(slide, Inches(6.8), y + Inches(0.02), Inches(1.8), Inches(0.35),
            title, sz=15, color=color, bold=True)
        _tb(slide, Inches(8.7), y + Inches(0.04), Inches(3.5), Inches(0.35),
            desc, sz=13, color=SUB_TEXT)

    # 底部：产品定位一句话
    _rrect(slide, Inches(1.16), Inches(6.3), Inches(11.0), Inches(0.5),
           fill=RGBColor(0xFD, 0xF0, 0xF0), border=SCU_RED)
    _tb(slide, Inches(1.16), Inches(6.33), Inches(11.0), Inches(0.45),
        "核心定位：AI 驱动的一站式智能校园助手  ·  对话+页面双模式  ·  拟对接真实教务系统",
        sz=14, color=SCU_RED, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# S08 — Part.03 项目概述与功能
# ============================================================
def s08_sec():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_bg(slide)
    _rect(slide, Inches(5.8), 0, Inches(7.5), SH, fill=RGBColor(0x2A, 0x0A, 0x14))
    _tb(slide, Inches(6.96), Inches(2.80), Inches(1.3), Inches(0.5),
        "Part.03", sz=18, color=GOLD, bold=True, font="Segoe UI")
    _tb(slide, Inches(5.68), Inches(3.31), Inches(6.1), Inches(1.0),
        "项目概述与功能", sz=36, color=WHITE, bold=True)
    _rect(slide, Inches(13.08), 0, Inches(0.25), SH, fill=SCU_RED)


# ============================================================
# S09 — 3.1 产品定位与交互模式
# ============================================================
def s09_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "3.1", "产品定位")
    _subtitle(slide, Inches(1.16), Inches(1.20), "AI 驱动的一站式智能校园助手")

    # 双模式（类似模板S4的左右两列）
    modes = [
        ("对话模式", "自然语言交互，AI 理解意图并自动调用功能模块",
         "例：「帮我看看明天有什么课」→ 自动查询课表返回结果", SCU_RED),
        ("页面模式", "传统结构化页面，快速浏览与操作各功能模块",
         "课表周视图 / 食堂导航 / DDL 看板 / 校车时刻", ACCENT_BLUE),
    ]
    for i, (title, desc, example, color) in enumerate(modes):
        x = Inches(1.16) + Inches(i * 5.8)
        y = Inches(1.80)
        _rrect(slide, x, y, Inches(5.5), Inches(2.0), fill=NEAR_WHITE, border=color)
        _rect(slide, x, y, Inches(5.5), Inches(0.5), fill=color)
        _tb(slide, x + Inches(0.2), y + Inches(0.07), Inches(5.0), Inches(0.4),
            title, sz=18, color=WHITE, bold=True)
        _tb(slide, x + Inches(0.3), y + Inches(0.6), Inches(4.9), Inches(0.5),
            desc, sz=14, color=DARK_TEXT, line_spacing=1.4)
        _tb(slide, x + Inches(0.3), y + Inches(1.3), Inches(4.9), Inches(0.4),
            example, sz=12, color=SUB_TEXT, italic=True)

    # 核心亮点（模板底部横排卡片）
    _tb(slide, Inches(0.91), Inches(4.1), Inches(3), Inches(0.4),
        "核心亮点", sz=18, color=DARK_TEXT, bold=True)

    highlights = [
        ("真实数据", "拟直连川大教务\nCAS + 验证码", SCU_RED),
        ("AI 原生", "拟接入通义千问\nFunction Calling", PURPLE),
        ("现代全栈", "Next.js + FastAPI\n+ Docker", ACCENT_BLUE),
        ("体验优先", "响应式设计\n+ 移动端适配", GREEN),
        ("工程化", "Git Flow 协作\n+ Docker 部署", ORANGE),
    ]
    for i, (title, desc, color) in enumerate(highlights):
        x = Inches(0.91) + Inches(i * 2.35)
        y = Inches(4.6)
        _rrect(slide, x, y, Inches(2.2), Inches(2.0), fill=WHITE, border=color)
        # 顶部色条
        _rect(slide, x, y, Inches(2.2), Inches(0.06), fill=color)
        # 图标圆
        _circle(slide, x + Inches(0.7), y + Inches(0.2), Inches(0.7), color)
        _tb(slide, x + Inches(0.72), y + Inches(0.28), Inches(0.7), Inches(0.55),
            title[0], sz=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 标题
        _tb(slide, x + Inches(0.1), y + Inches(1.0), Inches(2.0), Inches(0.3),
            title, sz=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        # 描述
        _tb(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.0), Inches(0.6),
            desc, sz=11, color=SUB_TEXT, align=PP_ALIGN.CENTER)


# ============================================================
# S10 — 3.2 核心功能模块
# ============================================================
def s10_features():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "3.2", "核心功能模块")

    modules = [
        ("AI 智能对话", SCU_RED, [
            "自然语言交互 + 多轮上下文",
            "Function Calling 意图路由",
            "个性化记忆系统",
        ]),
        ("课表 & 成绩", ACCENT_BLUE, [
            "周视图课程表 + 颜色区分",
            "成绩查询 & GPA 计算",
            "对接真实教务系统",
        ]),
        ("食堂导航", ORANGE, [
            "6 大食堂实时状态",
            "窗口导览与分类",
            "AI 推荐今天吃什么",
        ]),
        ("DDL 管理", GREEN, [
            "作业/考试截止日期",
            "优先级排序 + 提醒",
            "日历视图 + 逾期标记",
        ]),
        ("校车时刻", PURPLE, [
            "望江/江安/华西线路",
            "工作日/周末时刻表",
            "下一班倒计时",
        ]),
        ("更多（规划中）", MID_GRAY, [
            "校历查询 + 选课推荐",
            "RAG 文档问答",
            "晨间简报 + 深色模式",
        ]),
    ]

    # 3x2 网格
    col_w = Inches(3.65)
    row_h = Inches(2.35)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    start_x = Inches(0.91)
    start_y = Inches(1.20)

    for idx, (title, color, items) in enumerate(modules):
        col = idx % 3
        row = idx // 3
        x = start_x + (col_w + gap_x) * col
        y = start_y + (row_h + gap_y) * row

        _rrect(slide, x, y, col_w, row_h, fill=WHITE, border=color)
        _rect(slide, x, y, col_w, Inches(0.5), fill=color)
        _tb(slide, x + Inches(0.2), y + Inches(0.07), col_w - Inches(0.4), Inches(0.4),
            title, sz=16, color=WHITE, bold=True)
        for j, item in enumerate(items):
            _tb(slide, x + Inches(0.25), y + Inches(0.65) + Inches(j * 0.52),
                col_w - Inches(0.5), Inches(0.4),
                f"▸  {item}", sz=13, color=DARK_TEXT)


# ============================================================
# S11 — Part.04 团队成员与分工
# ============================================================
def s11_sec():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_bg(slide)
    _rect(slide, Inches(5.8), 0, Inches(7.5), SH, fill=RGBColor(0x2A, 0x0A, 0x14))
    _tb(slide, Inches(6.96), Inches(2.80), Inches(1.3), Inches(0.5),
        "Part.04", sz=18, color=GOLD, bold=True, font="Segoe UI")
    _tb(slide, Inches(5.68), Inches(3.31), Inches(6.1), Inches(1.0),
        "团队成员与分工", sz=36, color=WHITE, bold=True)
    _rect(slide, Inches(13.08), 0, Inches(0.25), SH, fill=SCU_RED)


# ============================================================
# S12 — 4.1 团队组织
# ============================================================
def s12_team():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "4.1", "团队组织")

    _tb(slide, Inches(4), Inches(0.60), Inches(6), Inches(0.4),
        "9 人团队  ·  4 个专业小组  ·  敏捷协作", sz=14, color=SUB_TEXT)

    teams = [
        ("前端组（2人）", ACCENT_BLUE, [
            ("谭博文", "Next.js 页面开发"),
            ("徐锐学", "UI/UX 组件搭建"),
        ]),
        ("后端组（3人）", GREEN, [
            ("孔垂骄", "FastAPI 接口"),
            ("张傲楚", "数据库建模"),
            ("毛立业", "教务系统对接"),
        ]),
        ("AI 组（2人）", PURPLE, [
            ("朱圣相", "LLM 意图路由"),
            ("李亚飞", "Function Calling"),
        ]),
        ("运维与管理（2人）", ORANGE, [
            ("谭旭睿", "Docker 部署"),
            ("覃泽锴", "项目管理"),
        ]),
    ]

    card_w = Inches(2.7)
    gap = Inches(0.2)
    start_x = Inches(0.91)
    start_y = Inches(1.20)

    for i, (group, color, members) in enumerate(teams):
        x = start_x + (card_w + gap) * i
        _rrect(slide, x, start_y, card_w, Inches(4.0), fill=WHITE, border=color)
        _rect(slide, x, start_y, card_w, Inches(0.5), fill=color)
        _tb(slide, x + Inches(0.1), start_y + Inches(0.07), card_w - Inches(0.2),
            Inches(0.4), group, sz=14, color=WHITE, bold=True)

        for j, (name, role) in enumerate(members):
            my = start_y + Inches(0.7) + Inches(j * 1.0)
            if name:
                _circle(slide, x + Inches(0.15), my + Inches(0.05), Inches(0.4), color)
                _tb(slide, x + Inches(0.17), my + Inches(0.1), Inches(0.4), Inches(0.35),
                    name[0], sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
                _tb(slide, x + Inches(0.65), my, Inches(1.8), Inches(0.3),
                    name, sz=14, color=DARK_TEXT, bold=True)
            _tb(slide, x + Inches(0.65), my + Inches(0.32), Inches(1.8), Inches(0.3),
                role, sz=12, color=SUB_TEXT)

    # 协作模式
    _rrect(slide, Inches(0.91), Inches(5.5), Inches(11.0), Inches(0.5),
           fill=RGBColor(0xFD, 0xF0, 0xF0), border=SCU_RED)
    _tb(slide, Inches(0.91), Inches(5.53), Inches(11.0), Inches(0.45),
        "协作模式：GitHub Flow  |  每周例会  |  双周迭代  |  接口文档先行",
        sz=13, color=SCU_RED, bold=True, align=PP_ALIGN.CENTER)

    # 工作量分配条
    bar_y = Inches(6.15)
    total_w = Inches(11.0)
    parts = [
        ("前端 25%", 0.25, ACCENT_BLUE),
        ("后端 30%", 0.30, GREEN),
        ("AI 25%",   0.25, PURPLE),
        ("运维 20%",  0.20, ORANGE),
    ]
    bx = Inches(0.91)
    for label, pct, color in parts:
        pw = int(total_w * pct)
        _rrect(slide, bx, bar_y, pw, Inches(0.4), fill=color)
        _tb(slide, bx, bar_y + Inches(0.04), pw, Inches(0.35),
            label, sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        bx += pw


# ============================================================
# S13 — Part.05 技术路线
# ============================================================
def s13_sec():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_bg(slide)
    _rect(slide, Inches(5.8), 0, Inches(7.5), SH, fill=RGBColor(0x2A, 0x0A, 0x14))
    _tb(slide, Inches(6.96), Inches(2.80), Inches(1.3), Inches(0.5),
        "Part.05", sz=18, color=GOLD, bold=True, font="Segoe UI")
    _tb(slide, Inches(5.68), Inches(3.31), Inches(6.1), Inches(1.0),
        "技术路线", sz=36, color=WHITE, bold=True)
    _rect(slide, Inches(13.08), 0, Inches(0.25), SH, fill=SCU_RED)


# ============================================================
# S14 — 5.1 系统架构
# ============================================================
def s14_architecture():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "5.1", "系统架构")

    # 三层架构（使用模板S15样式：三个横向条）
    layers = [
        ("前端展示层", "Next.js 14 + React 19 + shadcn/ui + TailwindCSS + Zustand",
         ACCENT_BLUE, RGBColor(0xEF, 0xF6, 0xFF)),
        ("后端服务层", "FastAPI + SQLAlchemy + Redis + JWT + Pydantic + httpx",
         GREEN, RGBColor(0xEC, 0xFD, 0xF5)),
        ("数据存储层", "PostgreSQL + Redis Cache + 教务系统 API + 通义千问 LLM",
         GOLD, RGBColor(0xFE, 0xF9, 0xEF)),
    ]

    layer_w = Inches(11.0)
    lx = Inches(0.91)

    for i, (name, techs, color, bg) in enumerate(layers):
        y = Inches(1.30) + Inches(i * 1.7)
        _rrect(slide, lx, y, layer_w, Inches(1.2), fill=bg, border=color)
        # 左侧标签
        _rect(slide, lx, y, Inches(0.08), Inches(1.2), fill=color)
        _tb(slide, lx + Inches(0.25), y + Inches(0.1), Inches(1.8), Inches(0.4),
            name, sz=16, color=color, bold=True)
        # 技术标签
        tech_list = techs.split(" + ")
        for j, tech in enumerate(tech_list):
            tx = lx + Inches(2.2) + Inches(j * 1.4)
            _rrect(slide, tx, y + Inches(0.5), Inches(1.25), Inches(0.5),
                   fill=WHITE, border=color)
            _tb(slide, tx, y + Inches(0.55), Inches(1.25), Inches(0.4),
                tech.strip(), sz=11, color=color, bold=True, align=PP_ALIGN.CENTER)

        # 层间箭头
        if i < 2:
            arrow_y = y + Inches(1.2)
            _tb(slide, Inches(5.5), arrow_y + Inches(0.05), Inches(2.5), Inches(0.35),
                "▼  REST API / SSE  ▼" if i == 0 else "▼  SQL / Cache / HTTP  ▼",
                sz=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # 选型原则
    _rrect(slide, lx, Inches(6.35), layer_w, Inches(0.5),
           fill=RGBColor(0xFD, 0xF0, 0xF0), border=SCU_RED)
    _tb(slide, lx, Inches(6.38), layer_w, Inches(0.45),
        "选型原则：  生产级成熟度  |  社区活跃度  |  团队学习曲线  |  异步优先性能",
        sz=14, color=SCU_RED, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# S15 — Part.06 展示结果与未来计划
# ============================================================
def s15_sec():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_bg(slide)
    _rect(slide, Inches(5.8), 0, Inches(7.5), SH, fill=RGBColor(0x2A, 0x0A, 0x14))
    _tb(slide, Inches(6.96), Inches(2.80), Inches(1.3), Inches(0.5),
        "Part.06", sz=18, color=GOLD, bold=True, font="Segoe UI")
    _tb(slide, Inches(5.68), Inches(3.31), Inches(6.1), Inches(1.0),
        "展示结果与未来计划", sz=36, color=WHITE, bold=True)
    _rect(slide, Inches(13.08), 0, Inches(0.25), SH, fill=SCU_RED)


# ============================================================
# S16 — 6.1 当前进展
# ============================================================
def s16_progress():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "6.1", "当前进展")

    # 总进度条
    bar_x = Inches(0.91)
    bar_w = Inches(11.0)
    bar_y = Inches(1.20)
    _rrect(slide, bar_x, bar_y, bar_w, Inches(0.45),
           fill=RGBColor(0xEE, 0xEE, 0xEE))
    _rrect(slide, bar_x, bar_y, int(bar_w * 0.10), Inches(0.5), fill=ACCENT_BLUE)
    _tb(slide, bar_x, bar_y + Inches(0.05), int(bar_w * 0.45), Inches(0.4),
        "10% — 规划阶段，正在学习技术栈", sz=15, color=DARK_TEXT, bold=True,
        align=PP_ALIGN.CENTER)

    # 已完成
    _rrect(slide, Inches(0.91), Inches(1.90), Inches(5.3), Inches(4.7),
           fill=WHITE, border=GREEN)
    _rect(slide, Inches(0.91), Inches(1.90), Inches(5.3), Inches(0.5), fill=GREEN)
    _tb(slide, Inches(1.1), Inches(1.95), Inches(3), Inches(0.4),
        "已完成", sz=15, color=WHITE, bold=True)

    done = [
        "确定项目选题与整体方案",
        "团队人员分工明确",
        "技术路线与技术栈选型完成",
        "项目功能模块规划完成",
    ]
    for i, item in enumerate(done):
        _tb(slide, Inches(1.2), Inches(2.55) + Inches(i * 0.55),
            Inches(4.8), Inches(0.4),
            f"✓  {item}", sz=13, color=DARK_TEXT)

    # 进行中 / 待完成
    _rrect(slide, Inches(6.5), Inches(1.90), Inches(5.3), Inches(4.7),
           fill=WHITE, border=ORANGE)
    _rect(slide, Inches(6.5), Inches(1.90), Inches(5.3), Inches(0.5), fill=ORANGE)
    _tb(slide, Inches(6.7), Inches(1.95), Inches(4), Inches(0.4),
        "进行中 / 待完成", sz=15, color=WHITE, bold=True)

    todo = [
        "学习 Next.js / React 前端技术",
        "学习 FastAPI / Python 后端技术",
        "学习 AI 大模型接口调用方式",
        "搭建前后端项目基础框架",
        "数据库表结构设计",
        "教务系统登录流程调研",
        "各功能模块开发（课表/食堂等）",
        "Docker 部署环境搭建",
    ]
    for i, item in enumerate(todo):
        _tb(slide, Inches(6.8), Inches(2.55) + Inches(i * 0.50),
            Inches(4.8), Inches(0.4),
            f"○  {item}", sz=13, color=DARK_TEXT)


# ============================================================
# S17 — 6.2 未来发展方向（模板S28风格）
# ============================================================
def s17_future():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _add_content_bg(slide)
    _title_with_line(slide, Inches(0.91), Inches(0.50), "6.2", "开发计划（12周）")

    # 四个阶段，按周排列
    phases = [
        ("第1-3周", "需求分析与框架搭建", ACCENT_BLUE,
         "需求调研 + 技术选型 + 前后端项目初始化 + 数据库设计 + 基础路由与登录"),
        ("第4-6周", "核心功能开发", GREEN,
         "课表/成绩页面 + 食堂导航 + 校车时刻 + DDL管理 + 教务CAS登录对接"),
        ("第7-9周", "AI接入与功能完善", PURPLE,
         "通义千问对话接入 + Function Calling + Dashboard汇总 + 深色模式"),
        ("第10-12周", "测试优化与交付", ORANGE,
         "性能优化 + 安全加固 + 用户测试 + Docker部署 + 项目文档 + 答辩准备"),
    ]

    # 时间轴横线
    line_y = Inches(1.50)
    _rect(slide, Inches(0.91), line_y, Inches(11.0), Pt(3), fill=SCU_RED)

    # 四个节点圆点
    for i in range(4):
        cx = Inches(0.91) + Inches(i * 3.1) + Inches(1.25)
        _circle(slide, cx, line_y - Inches(0.12), Inches(0.25), SCU_RED)

    # 四个阶段卡片（上下交错）
    for i, (weeks, title, color, desc) in enumerate(phases):
        x = Inches(0.50) + Inches(i * 2.95)
        y = Inches(2.0)
        card_w = Inches(2.85)
        card_h = Inches(3.8)

        _rrect(slide, x, y, card_w, card_h, fill=WHITE, border=color)
        # 顶部色条
        _rect(slide, x, y, card_w, Inches(0.45), fill=color)
        _tb(slide, x + Inches(0.1), y + Inches(0.05), card_w - Inches(0.2), Inches(0.35),
            weeks, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 标题
        _tb(slide, x + Inches(0.15), y + Inches(0.55), card_w - Inches(0.3), Inches(0.4),
            title, sz=14, color=color, bold=True, align=PP_ALIGN.CENTER)
        # 内容（逐条）
        items = desc.split(" + ")
        for j, item in enumerate(items):
            _tb(slide, x + Inches(0.2), y + Inches(1.05) + Inches(j * 0.48),
                card_w - Inches(0.4), Inches(0.4),
                f"▸ {item}", sz=12, color=DARK_TEXT)

    # 预期成果
    _rrect(slide, Inches(0.91), Inches(6.10), Inches(11.0), Inches(0.55),
           fill=RGBColor(0xFD, 0xF0, 0xF0), border=SCU_RED)
    _tb(slide, Inches(0.91), Inches(6.15), Inches(11.0), Inches(0.45),
        "预期成果：可运行的Web应用  |  AI基础问答能力  |  核心功能模块完整  |  Docker可部署  |  项目文档",
        sz=14, color=SCU_RED, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# S18 — 致谢页（复用模板封面背景）
# ============================================================
def s18_thanks():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(BG_COVER, Inches(-0.02), 0,
                             Inches(13.35), Inches(7.5))
    _add_overlay(slide, RGBColor(0x00, 0x00, 0x00), '30000')  # 30%黑色遮罩
    slide.shapes.add_picture(LOGO_BIG, Inches(0.15), Inches(0.15),
                             Inches(2.4), Inches(0.7))
    # 感谢文字（模板S30位置: 2.34, 2.85）
    _tb(slide, Inches(2.34), Inches(2.85), Inches(8.6), Inches(1.1),
        "感谢各位老师指正", sz=44, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)


# ============================================================
# 生成
# ============================================================
print("生成答辩风 PPT v3（基于模板）...")

s01_cover()
s02_toc()
s03_sec()           # Part.01 问题与出发点
s04_background()    # 1.1 选题背景
s05_data()          # 1.2 数据支撑
s06_sec()           # Part.02 需求分析
s07_requirement()   # 2.1 目标用户与核心需求
s08_sec()           # Part.03 项目概述与功能
s09_overview()      # 3.1 产品定位
s10_features()      # 3.2 核心功能模块
s11_sec()           # Part.04 团队成员与分工
s12_team()          # 4.1 团队组织
s13_sec()           # Part.05 技术路线
s14_architecture()  # 5.1 系统架构
s15_sec()           # Part.06 展示结果与未来计划
s16_progress()      # 6.1 当前进展
s17_future()        # 6.2 未来发展方向
s18_thanks()        # 致谢

out = os.path.join(SCRIPT_DIR, "SCU_Assistant_开题报告_v3k.pptx")
prs.save(out)
print(f"✓ 已生成: {out}")
print(f"  共 {len(prs.slides)} 页")
