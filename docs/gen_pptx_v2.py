"""
SCU Assistant 开题报告 PPT v2 — 基于 LaTeX Beamer 第二版内容
腾讯蓝渐变风格，输出 .pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 品牌色
# ============================================================
TENC_DARK  = RGBColor(0x00, 0x27, 0x66)
TENC_BLUE  = RGBColor(0x00, 0x52, 0xD9)
TENC_LIGHT = RGBColor(0x3B, 0x82, 0xF6)
TENC_BG    = RGBColor(0xEB, 0xF2, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG    = RGBColor(0xF8, 0xFA, 0xFF)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
SUB_TEXT   = RGBColor(0x64, 0x74, 0x8B)
SUCCESS    = RGBColor(0x00, 0xB5, 0x78)
WARN       = RGBColor(0xFF, 0x88, 0x00)
DANGER     = RGBColor(0xE3, 0x4D, 0x59)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
GOLD       = RGBColor(0xD9, 0x77, 0x06)
SCU_RED    = RGBColor(0xC4, 0x12, 0x30)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ============================================================
# 工具函数
# ============================================================

def _rect(slide, left, top, w, h, fill=None, border=None, border_w=Pt(0.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    return s

def _rrect(slide, left, top, w, h, fill=None, border=None, border_w=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    return s

def _circle(slide, left, top, d, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def _tb(slide, left, top, w, h, text, sz=14, color=DARK_TEXT,
        bold=False, italic=False, font="Microsoft YaHei",
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.15):
    """添加文本框"""
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    p.line_spacing = Pt(sz * line_spacing)
    tf.auto_size = None
    return tf

def _multi_tb(slide, left, top, w, h, lines, default_sz=12, default_color=DARK_TEXT,
              align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    """多行文本框，lines = [(text, sz, color, bold), ...]"""
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(lines):
        text = item[0]
        sz = item[1] if len(item) > 1 else default_sz
        color = item[2] if len(item) > 2 else default_color
        bold = item[3] if len(item) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.alignment = align
        p.line_spacing = Pt(sz * 1.4)
    return tf


def _page_bg(slide):
    """标准内容页背景：白底 + 蓝色顶栏"""
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    # 蓝色顶栏
    _rect(slide, 0, 0, SW, Inches(0.95), fill=TENC_BLUE)
    # 底部灰线
    _rect(slide, 0, Inches(7.15), SW, Inches(0.35), fill=RGBColor(0xF1, 0xF5, 0xF9))

def _page_title(slide, title, subtitle=""):
    """顶栏标题"""
    _tb(slide, Inches(0.6), Inches(0.18), Inches(8), Inches(0.55),
        title, sz=22, color=WHITE, bold=True)
    if subtitle:
        _tb(slide, Inches(0.6), Inches(0.55), Inches(8), Inches(0.35),
            subtitle, sz=11, color=RGBColor(0xBB, 0xD5, 0xFF))

def _page_number(slide, num, total):
    _tb(slide, Inches(12.3), Inches(7.18), Inches(0.8), Inches(0.3),
        f"{num}/{total}", sz=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def _card(slide, left, top, w, h, fill=WHITE, border=None):
    """白色圆角卡片"""
    return _rrect(slide, left, top, w, h, fill=fill,
                  border=border or RGBColor(0x00, 0x52, 0xD9),
                  border_w=Pt(0.5))

def _blue_card_title(slide, left, top, w, title, title_color=TENC_BLUE):
    """卡片标题（蓝色加粗）"""
    _tb(slide, left, top, w, Inches(0.35), title,
        sz=13, color=title_color, bold=True)

def _section_slide(slide, part_num, title):
    """章节分隔页"""
    _rect(slide, 0, 0, SW, SH, fill=TENC_DARK)
    _circle(slide, Inches(8.5), Inches(-1.5), Inches(8), RGBColor(0x00, 0x3D, 0x99))
    _tb(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
        f"PART {part_num}", sz=18, color=MID_GRAY, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1), Inches(3.3), Inches(11), Inches(1),
        title, sz=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


TOTAL_SLIDES = 15  # 含 section 分隔页


# ============================================================
# S01 — 封面
# ============================================================
def s01_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=TENC_DARK)
    _rect(slide, 0, 0, Inches(6.5), SH, fill=TENC_BLUE)
    _circle(slide, Inches(10), Inches(-0.5), Inches(4), RGBColor(0x1A, 0x3D, 0x7A))
    _circle(slide, Inches(8.5), Inches(4.5), Inches(5), RGBColor(0x10, 0x30, 0x60))
    _rect(slide, 0, Inches(7.3), SW, Inches(0.03), fill=RGBColor(0x30, 0x30, 0x60))

    _tb(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
        "SCU Assistant", sz=48, color=WHITE, bold=True, font="Segoe UI")
    _tb(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.7),
        "四川大学智能校园助手", sz=30, color=RGBColor(0xBB, 0xD5, 0xFF), bold=True)
    _rect(slide, Inches(1.5), Inches(3.8), Inches(2.5), Inches(0.03), fill=TENC_LIGHT)
    _tb(slide, Inches(1.5), Inches(4.1), Inches(8), Inches(0.4),
        "项目开题报告", sz=18, color=RGBColor(0xDD, 0xDD, 0xFF))
    _tb(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.4),
        "四川大学计算机学院  |  软件工程课程设计  |  2026 春季学期",
        sz=12, color=MID_GRAY)


# ============================================================
# S02 — 目录
# ============================================================
def s02_toc():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "目录", "CONTENTS")
    _page_number(slide, 2, TOTAL_SLIDES)

    items_left = [
        ("01", "项目背景与痛点"),
        ("02", "产品定位与目标"),
        ("03", "核心功能模块"),
        ("04", "系统架构设计"),
        ("05", "技术选型"),
    ]
    items_right = [
        ("06", "团队成员与分工"),
        ("07", "开发计划与里程碑"),
        ("08", "技术难点与风险"),
        ("09", "当前进展"),
        ("10", "总结与展望"),
    ]

    # 左卡片
    _card(slide, Inches(1.2), Inches(1.4), Inches(5.2), Inches(5.0))
    for i, (num, title) in enumerate(items_left):
        y = Inches(1.7) + Inches(i * 0.85)
        _tb(slide, Inches(1.6), y, Inches(0.6), Inches(0.4),
            num, sz=20, color=TENC_BLUE, bold=True)
        _tb(slide, Inches(2.4), y + Inches(0.05), Inches(3.5), Inches(0.35),
            title, sz=15, color=DARK_TEXT)

    # 右卡片
    _card(slide, Inches(7.0), Inches(1.4), Inches(5.2), Inches(5.0))
    for i, (num, title) in enumerate(items_right):
        y = Inches(1.7) + Inches(i * 0.85)
        _tb(slide, Inches(7.4), y, Inches(0.6), Inches(0.4),
            num, sz=20, color=TENC_BLUE, bold=True)
        _tb(slide, Inches(8.2), y + Inches(0.05), Inches(3.5), Inches(0.35),
            title, sz=15, color=DARK_TEXT)


# ============================================================
# S03 — Section: 项目背景
# ============================================================
def s03_sec_bg():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _section_slide(slide, "01", "项目背景与痛点")


# ============================================================
# S04 — 项目背景
# ============================================================
def s04_background():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "项目背景", "校园信息化现状")
    _page_number(slide, 4, TOTAL_SLIDES)

    # 左侧文字
    _tb(slide, Inches(0.7), Inches(1.2), Inches(5.8), Inches(0.8),
        "四川大学在校学生超过 5 万人，日常需频繁访问多个独立平台获取校园信息。",
        sz=13, color=DARK_TEXT)

    # 四大核心痛点卡片
    _card(slide, Inches(0.7), Inches(2.0), Inches(5.8), Inches(3.2),
          border=RGBColor(0x00, 0x52, 0xD9))
    _blue_card_title(slide, Inches(1.0), Inches(2.15), Inches(4), "四大核心痛点")

    pains = [
        ("信息碎片化", "课表、校车、食堂分布在 5+ 个平台"),
        ("系统体验差", "教务 UI 落后，移动端适配差"),
        ("缺乏智能化", "无个性化推荐，无统一问答入口"),
        ("数据不互通", "各系统数据孤岛，时效性差"),
    ]
    for i, (t, d) in enumerate(pains):
        y = Inches(2.65) + Inches(i * 0.6)
        _tb(slide, Inches(1.3), y, Inches(1.5), Inches(0.3),
            f"▸ {t}", sz=12, color=TENC_BLUE, bold=True)
        _tb(slide, Inches(2.9), y, Inches(3.2), Inches(0.3),
            f"— {d}", sz=11, color=SUB_TEXT)

    # 右侧 — 信息聚合示意图
    _card(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.3))
    platforms = ["教务系统", "食堂公众号", "校车通知群", "教学日历", "各类QQ群"]
    for i, name in enumerate(platforms):
        y = Inches(1.6) + Inches(i * 0.65)
        _rrect(slide, Inches(8.5), y, Inches(2.5), Inches(0.5),
               fill=TENC_BG, border=RGBColor(0x00, 0x52, 0xD9))
        _tb(slide, Inches(8.6), y + Inches(0.05), Inches(2.3), Inches(0.4),
            name, sz=11, color=TENC_BLUE, align=PP_ALIGN.CENTER)
        # 箭头线
        if i < len(platforms):
            _rect(slide, Inches(9.7), y + Inches(0.5), Inches(0.04), Inches(0.15),
                  fill=MID_GRAY)

    # 汇聚目标
    _rrect(slide, Inches(8.0), Inches(5.3), Inches(3.5), Inches(0.6),
           fill=RGBColor(0xE8, 0xF0, 0xFE), border=TENC_BLUE)
    _tb(slide, Inches(8.1), Inches(5.35), Inches(3.3), Inches(0.5),
        "SCU Assistant（一站式）", sz=13, color=TENC_BLUE, bold=True,
        align=PP_ALIGN.CENTER)

    # 标签
    _tb(slide, Inches(11.3), Inches(3.0), Inches(1.2), Inches(0.3),
        "碎片化", sz=10, color=DANGER, bold=True)
    _tb(slide, Inches(11.5), Inches(5.5), Inches(1.2), Inches(0.3),
        "一站式", sz=10, color=SUCCESS, bold=True)


# ============================================================
# S05 — 数据支撑
# ============================================================
def s05_data():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "数据支撑")
    _page_number(slide, 5, TOTAL_SLIDES)

    stats = [
        ("50000+", "川大在校学生", TENC_BLUE),
        ("5+", "日常使用平台数", TENC_BLUE),
        ("70%", "反馈查课不便", DANGER),
        ("0", "现有智能助手", GOLD),
    ]

    card_w = Inches(2.5)
    gap = Inches(0.4)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) // 2

    for i, (num, label, color) in enumerate(stats):
        x = start_x + (card_w + gap) * i
        _card(slide, x, Inches(2.0), card_w, Inches(2.2))
        _tb(slide, x, Inches(2.3), card_w, Inches(0.9),
            num, sz=36, color=color, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, x, Inches(3.3), card_w, Inches(0.4),
            label, sz=13, color=SUB_TEXT, align=PP_ALIGN.CENTER)

    # 底部结论
    cw = Inches(10)
    cx = (SW - cw) // 2
    _card(slide, cx, Inches(5.0), cw, Inches(1.2))
    _tb(slide, cx, Inches(5.2), cw, Inches(0.4),
        "市场空白：川大尚无集成式智能校园助手",
        sz=16, color=TENC_BLUE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, cx, Inches(5.7), cw, Inches(0.4),
        "学生需要统一入口高效获取课表、食堂、校车、DDL 等信息",
        sz=12, color=SUB_TEXT, align=PP_ALIGN.CENTER)


# ============================================================
# S06 — Section: 产品定位
# ============================================================
def s06_sec_pos():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _section_slide(slide, "02", "产品定位与目标")


# ============================================================
# S07 — 产品定位
# ============================================================
def s07_positioning():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "产品定位")
    _page_number(slide, 7, TOTAL_SLIDES)

    # 核心定位卡片
    _card(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(1.3),
          fill=RGBColor(0xF0, 0xF5, 0xFF), border=TENC_BLUE)
    _blue_card_title(slide, Inches(1.0), Inches(1.4), Inches(4), "核心定位")
    _tb(slide, Inches(1.0), Inches(1.8), Inches(5.3), Inches(0.4),
        "AI 驱动的一站式智能校园助手", sz=16, color=TENC_BLUE, bold=True,
        align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1.0), Inches(2.2), Inches(5.3), Inches(0.3),
        "对话 + 页面双模式  |  对接真实教务  |  个性化记忆",
        sz=10, color=SUB_TEXT, align=PP_ALIGN.CENTER)

    # 目标用户
    _card(slide, Inches(0.7), Inches(2.9), Inches(5.8), Inches(2.5))
    _blue_card_title(slide, Inches(1.0), Inches(3.05), Inches(4), "目标用户")
    users = [
        "四川大学全日制本科生 & 研究生",
        "日均使用教务/食堂/校车的学生",
        "期望通过 AI 提升效率的用户",
    ]
    for i, u in enumerate(users):
        _tb(slide, Inches(1.3), Inches(3.55) + Inches(i * 0.5), Inches(4.8), Inches(0.4),
            f"▸  {u}", sz=12, color=DARK_TEXT)

    # 右侧 — 核心目标
    _tb(slide, Inches(7.0), Inches(1.3), Inches(3), Inches(0.4),
        "核心目标", sz=15, color=TENC_BLUE, bold=True)

    goals = [
        ("效率提升", "一句话查课表，减少80%步骤", TENC_BLUE),
        ("智能推荐", "AI意图识别 + 个性化", PURPLE),
        ("数据整合", "教务/食堂/校车一站式", GOLD),
        ("体验升级", "深色主题 + 移动端优先", SUCCESS),
    ]
    for i, (title, desc, color) in enumerate(goals):
        y = Inches(1.9) + Inches(i * 0.95)
        _rrect(slide, Inches(7.0), y, Inches(5.5), Inches(0.75),
               fill=WHITE, border=color)
        # 左色条
        _rect(slide, Inches(7.0), y, Inches(0.08), Inches(0.75), fill=color)
        _tb(slide, Inches(7.3), y + Inches(0.1), Inches(2), Inches(0.3),
            title, sz=13, color=color, bold=True)
        _tb(slide, Inches(9.3), y + Inches(0.12), Inches(3), Inches(0.3),
            desc, sz=10, color=SUB_TEXT)


# ============================================================
# S08 — Section: 核心功能
# ============================================================
def s08_sec_feat():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _section_slide(slide, "03", "核心功能模块")


# ============================================================
# S09 — P0 核心功能
# ============================================================
def s09_features():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "P0 核心功能")
    _page_number(slide, 9, TOTAL_SLIDES)

    modules = [
        ("AI 智能对话", SCU_RED, [
            "自然语言交互",
            "Function Calling 意图路由",
            "个性化记忆系统",
            "多轮上下文理解",
        ]),
        ("课表 & 成绩", TENC_BLUE, [
            "周视图课程表",
            "成绩查询 & GPA 计算",
            "学分进度追踪",
            "对接真实教务系统",
        ]),
        ("食堂导航", WARN, [
            "6 大食堂实时状态",
            "窗口导览与分类",
            "AI 推荐今天吃什么",
            "校区筛选",
        ]),
        ("DDL 管理", SUCCESS, [
            "作业/考试截止日期",
            "优先级排序",
            "逾期/紧急提醒",
        ]),
        ("校车时刻", PURPLE, [
            "望江/江安/华西线路",
            "工作日/周末时刻表",
            "下一班倒计时",
        ]),
        ("更多 (P1)", SUB_TEXT, [
            "校历查询",
            "选课推荐",
            "RAG 文档问答",
        ]),
    ]

    # 3 列 x 2 行
    col_w = Inches(3.7)
    row_h = Inches(2.6)
    gap_x = Inches(0.4)
    gap_y = Inches(0.3)
    start_x = Inches(0.7)
    start_y = Inches(1.2)

    for idx, (title, color, items) in enumerate(modules):
        col = idx % 3
        row = idx // 3
        x = start_x + (col_w + gap_x) * col
        y = start_y + (row_h + gap_y) * row

        # 卡片
        _rrect(slide, x, y, col_w, row_h, fill=WHITE, border=color)
        # 标题栏
        _rect(slide, x, y, col_w, Inches(0.5), fill=color)
        _tb(slide, x + Inches(0.2), y + Inches(0.07), col_w - Inches(0.4), Inches(0.4),
            title, sz=13, color=WHITE, bold=True)
        # 列表
        for j, item in enumerate(items):
            _tb(slide, x + Inches(0.3), y + Inches(0.65) + Inches(j * 0.42),
                col_w - Inches(0.6), Inches(0.35),
                f"▸  {item}", sz=11, color=DARK_TEXT)


# ============================================================
# S10 — Section: 系统架构
# ============================================================
def s10_sec_arch():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _section_slide(slide, "04", "系统架构设计")


# ============================================================
# S11 — 整体架构
# ============================================================
def s11_architecture():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "整体架构", "前后端分离 + 模块化服务")
    _page_number(slide, 11, TOTAL_SLIDES)

    layer_w = Inches(11.5)
    layer_h = Inches(1.3)
    lx = (SW - layer_w) // 2

    # ---- 前端展示层 ----
    y1 = Inches(1.3)
    _rrect(slide, lx, y1, layer_w, layer_h, fill=RGBColor(0xEF, 0xF6, 0xFF), border=TENC_BLUE)
    _tb(slide, lx + Inches(0.2), y1 + Inches(0.1), Inches(2), Inches(0.3),
        "前端展示层", sz=12, color=TENC_BLUE, bold=True)
    _tb(slide, lx + Inches(0.2), y1 + Inches(0.45), Inches(4), Inches(0.3),
        "Next.js 14 + React 19 + shadcn/ui", sz=9, color=SUB_TEXT)

    fe_modules = ["登录页", "Dashboard", "课表", "成绩", "食堂", "校车", "AI 对话"]
    fe_colors = [TENC_BLUE]*6 + [SCU_RED]
    for i, (name, c) in enumerate(zip(fe_modules, fe_colors)):
        mx = lx + Inches(3.5) + Inches(i * 1.12)
        _rrect(slide, mx, y1 + Inches(0.35), Inches(0.95), Inches(0.55),
               fill=WHITE, border=c)
        _tb(slide, mx, y1 + Inches(0.4), Inches(0.95), Inches(0.45),
            name, sz=8, color=c, bold=True, align=PP_ALIGN.CENTER)

    # 箭头
    _tb(slide, Inches(6.4), y1 + layer_h, Inches(1), Inches(0.4),
        "REST / SSE", sz=8, color=SUB_TEXT, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(6.65), y1 + layer_h + Inches(0.05), Inches(0.04), Inches(0.35),
          fill=MID_GRAY)

    # ---- 后端服务层 ----
    y2 = y1 + layer_h + Inches(0.5)
    _rrect(slide, lx, y2, layer_w, layer_h, fill=RGBColor(0xEC, 0xFD, 0xF5), border=SUCCESS)
    _tb(slide, lx + Inches(0.2), y2 + Inches(0.1), Inches(2), Inches(0.3),
        "后端服务层", sz=12, color=SUCCESS, bold=True)
    _tb(slide, lx + Inches(0.2), y2 + Inches(0.45), Inches(4), Inches(0.3),
        "FastAPI + JWT + Redis", sz=9, color=SUB_TEXT)

    be_modules = [("Auth 认证", SUCCESS), ("Academic", SUCCESS), ("Food", SUCCESS),
                  ("Campus", SUCCESS), ("AI Intent", PURPLE), ("Rate Limit", WARN)]
    for i, (name, c) in enumerate(be_modules):
        mx = lx + Inches(3.5) + Inches(i * 1.3)
        _rrect(slide, mx, y2 + Inches(0.35), Inches(1.1), Inches(0.55),
               fill=WHITE, border=c)
        _tb(slide, mx, y2 + Inches(0.4), Inches(1.1), Inches(0.45),
            name, sz=8, color=c, bold=True, align=PP_ALIGN.CENTER)

    # 箭头
    _rect(slide, Inches(6.65), y2 + layer_h + Inches(0.05), Inches(0.04), Inches(0.35),
          fill=MID_GRAY)

    # ---- 数据存储层 ----
    y3 = y2 + layer_h + Inches(0.5)
    _rrect(slide, lx, y3, layer_w, layer_h, fill=RGBColor(0xFE, 0xF9, 0xEF), border=GOLD)
    _tb(slide, lx + Inches(0.2), y3 + Inches(0.1), Inches(2), Inches(0.3),
        "数据存储层", sz=12, color=GOLD, bold=True)
    _tb(slide, lx + Inches(0.2), y3 + Inches(0.45), Inches(5), Inches(0.3),
        "PostgreSQL + Redis + 外部API", sz=9, color=SUB_TEXT)

    db_modules = [("PostgreSQL", TENC_BLUE), ("Redis", WARN),
                  ("教务系统", SCU_RED), ("通义千问", PURPLE)]
    for i, (name, c) in enumerate(db_modules):
        mx = lx + Inches(3.5) + Inches(i * 1.8)
        _rrect(slide, mx, y3 + Inches(0.35), Inches(1.4), Inches(0.55),
               fill=WHITE, border=c)
        _tb(slide, mx, y3 + Inches(0.4), Inches(1.4), Inches(0.45),
            name, sz=9, color=c, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# S12 — 技术选型（无 section 分隔页，紧跟架构）
# ============================================================
def s12_techstack():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "技术选型详情")
    _page_number(slide, 12, TOTAL_SLIDES)

    groups = [
        ("前端", TENC_BLUE, [
            ("Next.js 14", "SSR 框架"),
            ("React 19", "UI 库"),
            ("TailwindCSS", "样式"),
            ("shadcn/ui", "组件库"),
            ("Zustand", "状态管理"),
            ("TanStack", "请求缓存"),
            ("Zod", "表单验证"),
        ]),
        ("后端", SUCCESS, [
            ("FastAPI", "Web 框架"),
            ("SQLAlchemy", "异步 ORM"),
            ("Alembic", "DB 迁移"),
            ("Redis 7", "缓存/限流"),
            ("PyJWT", "Token 认证"),
            ("Pydantic", "数据验证"),
            ("httpx", "HTTP 客户端"),
        ]),
        ("AI / 数据", PURPLE, [
            ("通义千问", "主力 LLM"),
            ("LangChain", "编排框架"),
            ("PostgreSQL", "主数据库"),
            ("JSONB", "灵活存储"),
            ("RAG", "文档问答"),
        ]),
        ("DevOps", WARN, [
            ("Docker", "容器化"),
            ("Compose", "编排部署"),
            ("Actions", "CI/CD"),
            ("Ruff", "Python Lint"),
            ("ESLint", "JS Lint"),
            ("pytest", "后端测试"),
            ("Vitest", "前端测试"),
        ]),
    ]

    card_w = Inches(2.8)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    start_y = Inches(1.2)
    card_h = Inches(4.8)

    for gi, (group_name, color, items) in enumerate(groups):
        x = start_x + (card_w + gap) * gi
        _rrect(slide, x, start_y, card_w, card_h, fill=WHITE, border=color)
        # 标题
        _rect(slide, x, start_y, card_w, Inches(0.5), fill=color)
        _tb(slide, x + Inches(0.2), start_y + Inches(0.07), card_w - Inches(0.4), Inches(0.4),
            group_name, sz=13, color=WHITE, bold=True)
        # 表格行
        for j, (tech, desc) in enumerate(items):
            ry = start_y + Inches(0.7) + Inches(j * 0.52)
            # 交替背景
            if j % 2 == 0:
                _rect(slide, x + Inches(0.05), ry, card_w - Inches(0.1), Inches(0.45),
                      fill=RGBColor(0xF8, 0xFA, 0xFF))
            _tb(slide, x + Inches(0.15), ry + Inches(0.05), Inches(1.2), Inches(0.35),
                tech, sz=10, color=DARK_TEXT, bold=True)
            _tb(slide, x + Inches(1.4), ry + Inches(0.05), Inches(1.2), Inches(0.35),
                desc, sz=10, color=SUB_TEXT)

    # 底部选型原则
    _card(slide, Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.65))
    _tb(slide, Inches(0.8), Inches(6.3), Inches(11.6), Inches(0.45),
        "选型原则：  生产级成熟度  |  社区活跃度  |  团队学习曲线  |  异步优先性能",
        sz=12, color=TENC_BLUE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# S13 — 团队成员与分工
# ============================================================
def s13_team():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "团队组织", "8 人团队 · 4 个专业小组 · 敏捷协作")
    _page_number(slide, 13, TOTAL_SLIDES)

    teams = [
        ("前端组 (2-3人)", TENC_BLUE, [
            "Next.js 页面开发",
            "UI/UX 组件搭建",
            "响应式适配",
            "状态管理",
            "API 对接",
        ]),
        ("后端组 (2-3人)", SUCCESS, [
            "FastAPI 接口",
            "数据库建模",
            "教务系统对接",
            "JWT 认证",
            "单元测试",
        ]),
        ("AI 组 (1-2人)", PURPLE, [
            "LLM 意图路由",
            "Function Calling",
            "用户记忆系统",
            "RAG 引擎",
            "Prompt 优化",
        ]),
        ("DevOps/PM (1人)", WARN, [
            "Docker 部署",
            "CI/CD 流水线",
            "项目管理",
            "Code Review",
            "文档与演示",
        ]),
    ]

    card_w = Inches(2.8)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    start_y = Inches(1.2)

    for i, (title, color, items) in enumerate(teams):
        x = start_x + (card_w + gap) * i
        _rrect(slide, x, start_y, card_w, Inches(3.8), fill=WHITE, border=color)
        _rect(slide, x, start_y, card_w, Inches(0.5), fill=color)
        _tb(slide, x + Inches(0.2), start_y + Inches(0.07), card_w - Inches(0.4), Inches(0.4),
            title, sz=12, color=WHITE, bold=True)
        for j, item in enumerate(items):
            _tb(slide, x + Inches(0.3), start_y + Inches(0.7) + Inches(j * 0.55),
                card_w - Inches(0.6), Inches(0.4),
                f"▸  {item}", sz=11, color=DARK_TEXT)

    # 底部协作模式
    _card(slide, Inches(0.6), Inches(5.3), Inches(12.0), Inches(0.7))
    _tb(slide, Inches(0.8), Inches(5.4), Inches(11.6), Inches(0.5),
        "协作模式：  GitHub Flow 分支管理  |  每周 Standup  |  双周 Sprint Review  |  OpenAPI 文档先行",
        sz=12, color=TENC_BLUE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# S14 — 16周开发路线图
# ============================================================
def s14_timeline():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "16 周开发路线图")
    _page_number(slide, 14, TOTAL_SLIDES)

    # 表格
    headers = ["Sprint", "阶段", "周", "核心任务", "交付物"]
    col_ws = [Inches(0.7), Inches(1.2), Inches(0.9), Inches(4.5), Inches(2.5)]
    start_x = Inches(0.8)
    start_y = Inches(1.2)
    row_h = Inches(0.55)

    # Header
    cx = start_x
    for hi, (hdr, cw) in enumerate(zip(headers, col_ws)):
        _rect(slide, cx, start_y, cw, row_h, fill=TENC_BLUE)
        _tb(slide, cx + Inches(0.05), start_y + Inches(0.08), cw - Inches(0.1), Inches(0.4),
            hdr, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cw

    rows = [
        ("0", "基础设施", "W1-2",  "Git、Docker Compose、CI/CD", "开发环境"),
        ("1", "核心骨架", "W3-4",  "登录认证、主布局、DB迁移", "可登录原型"),
        ("2", "学业MVP",  "W5-6",  "课表+成绩+教务系统对接", "学业模块"),
        ("3", "食堂+校车", "W7-8", "食堂导航、校车时刻、DDL", "生活模块"),
        ("4", "AI 对话",  "W9-10", "意图路由、Function Calling", "AI 模块"),
        ("5", "P1 扩展",  "W11-12","选课推荐、RAG 文档问答", "增强功能"),
        ("6", "打磨优化", "W13-14","晨间简报、UI精调、深色模式", "体验升级"),
        ("7", "测试交付", "W15-16","Bug修复、文档、演示准备", "最终交付"),
    ]

    for ri, row_data in enumerate(rows):
        ry = start_y + row_h * (ri + 1)
        bg = TENC_BG if ri % 2 == 1 else WHITE
        cx = start_x
        for ci, (val, cw) in enumerate(zip(row_data, col_ws)):
            _rect(slide, cx, ry, cw, row_h, fill=bg)
            color = TENC_BLUE if ci == 0 else DARK_TEXT
            bold = ci <= 1
            _tb(slide, cx + Inches(0.05), ry + Inches(0.08), cw - Inches(0.1), Inches(0.4),
                val, sz=10, color=color, bold=bold, align=PP_ALIGN.CENTER if ci <= 2 else PP_ALIGN.LEFT)
            cx += cw

    # 时间轴
    tl_y = Inches(6.0)
    tl_x = Inches(1.2)
    tl_w = Inches(11)
    _rect(slide, tl_x, tl_y + Inches(0.15), tl_w, Inches(0.04), fill=MID_GRAY)

    milestones = [
        (0.0, "W1 启动", TENC_BLUE),
        (0.25, "W4 骨架验收", GOLD),
        (0.5, "W8 MVP", SUCCESS),
        (0.625, "W10 AI上线", PURPLE),
        (1.0, "W16 最终交付", SCU_RED),
    ]
    for pct, label, color in milestones:
        mx = tl_x + int(tl_w * pct)
        _circle(slide, mx - Inches(0.12), tl_y, Inches(0.3), color)
        _tb(slide, mx - Inches(0.5), tl_y + Inches(0.35), Inches(1.2), Inches(0.3),
            label, sz=8, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # 当前标记
    cur_x = tl_x + int(tl_w * 0.35)
    _rect(slide, cur_x, tl_y - Inches(0.15), Inches(0.03), Inches(0.6), fill=SCU_RED)
    _rrect(slide, cur_x - Inches(0.3), tl_y - Inches(0.35), Inches(0.65), Inches(0.25),
           fill=SCU_RED)
    _tb(slide, cur_x - Inches(0.3), tl_y - Inches(0.33), Inches(0.65), Inches(0.22),
        "当前 W6", sz=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# S15 — 风险分析
# ============================================================
def s15_risks():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "风险分析与应对")
    _page_number(slide, 15, TOTAL_SLIDES)

    # 表格
    headers = ["风险项", "等级", "挑战", "应对方案"]
    col_ws = [Inches(2.0), Inches(0.8), Inches(4.2), Inches(4.8)]
    start_x = Inches(0.6)
    start_y = Inches(1.2)
    row_h = Inches(0.7)

    cx = start_x
    for hdr, cw in zip(headers, col_ws):
        _rect(slide, cx, start_y, cw, Inches(0.5), fill=TENC_BLUE)
        _tb(slide, cx + Inches(0.05), start_y + Inches(0.05), cw - Inches(0.1), Inches(0.4),
            hdr, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cw

    risks = [
        ("教务系统对接", "高", DANGER, "反爬 + 验证码 + Session管理", "双MD5加密还原；Mock并行开发"),
        ("AI 意图准确率", "中", WARN, "语言歧义 + 多意图混合", "Prompt优化；页面导航降级兜底"),
        ("性能与并发", "中", WARN, "高峰并发 + LLM响应延迟", "Redis多级缓存；SSE流式输出"),
        ("团队协作", "低", SUCCESS, "接口对齐 + 风格统一", "OpenAPI先行；CI自动Lint"),
    ]

    for ri, (name, level, level_color, challenge, solution) in enumerate(risks):
        ry = start_y + Inches(0.5) + row_h * ri
        bg = TENC_BG if ri % 2 == 0 else WHITE
        vals = [name, "", challenge, solution]
        cx = start_x
        for ci, (val, cw) in enumerate(zip(vals, col_ws)):
            _rect(slide, cx, ry, cw, row_h, fill=bg)
            if ci == 1:
                # 等级标签
                tag_w = Inches(0.5)
                tag_x = cx + (cw - tag_w) // 2
                tag_y = ry + Inches(0.15)
                _rrect(slide, tag_x, tag_y, tag_w, Inches(0.35), fill=level_color)
                _tb(slide, tag_x, tag_y + Inches(0.02), tag_w, Inches(0.3),
                    level, sz=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            else:
                _tb(slide, cx + Inches(0.1), ry + Inches(0.15), cw - Inches(0.2), Inches(0.4),
                    val, sz=10, color=DARK_TEXT, bold=(ci==0))
            cx += cw

    # 工作量分配条
    bar_y = Inches(4.8)
    _tb(slide, Inches(0.6), bar_y, Inches(3), Inches(0.35),
        "工作量分配", sz=13, color=TENC_BLUE, bold=True, align=PP_ALIGN.CENTER)

    bar_y2 = Inches(5.3)
    total_w = Inches(11.5)
    parts = [
        ("前端 35%", 0.35, TENC_BLUE),
        ("后端 30%", 0.30, SUCCESS),
        ("AI 20%",   0.20, PURPLE),
        ("Ops 15%",  0.15, WARN),
    ]
    bx = Inches(0.6)
    for label, pct, color in parts:
        pw = int(total_w * pct)
        _rect(slide, bx, bar_y2, pw, Inches(0.45), fill=color)
        _tb(slide, bx, bar_y2 + Inches(0.05), pw, Inches(0.35),
            label, sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        bx += pw


# ============================================================
# S16 — 当前进展
# ============================================================
def s16_progress():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "项目开发进度")
    _page_number(slide, 16, TOTAL_SLIDES)

    # 总进度条
    bar_x = Inches(1)
    bar_w = Inches(11)
    bar_y = Inches(1.3)
    _rrect(slide, bar_x, bar_y, bar_w, Inches(0.5), fill=TENC_BG)
    _rrect(slide, bar_x, bar_y, int(bar_w * 0.6), Inches(0.5), fill=SUCCESS)
    _tb(slide, bar_x, bar_y + Inches(0.05), int(bar_w * 0.6), Inches(0.4),
        "60% — P0 核心模块已完成", sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, bar_x + int(bar_w * 0.65), bar_y + Inches(0.08), Inches(2.5), Inches(0.35),
        "P1 功能开发中", sz=9, color=SUB_TEXT)

    # 已完成
    _rrect(slide, Inches(0.7), Inches(2.1), Inches(5.8), Inches(4.5),
           fill=WHITE, border=SUCCESS)
    _rect(slide, Inches(0.7), Inches(2.1), Inches(5.8), Inches(0.5), fill=SUCCESS)
    _tb(slide, Inches(0.9), Inches(2.17), Inches(3), Inches(0.4),
        "✓  已完成", sz=13, color=WHITE, bold=True)

    done_items = [
        "认证系统（教务真实登录 + JWT）",
        "Dashboard 仪表盘（今日课程/成绩/学分）",
        "课程表页面（周视图 + 颜色区分）",
        "DDL 管理（本地 CRUD + 优先级）",
        "食堂导航（6 食堂 + 窗口导览）",
        "校车时刻（三校区 + 倒计时）",
        "AI 对话 UI 骨架",
        "CI/CD 流水线 + Docker 部署",
    ]
    for i, item in enumerate(done_items):
        _tb(slide, Inches(1.1), Inches(2.75) + Inches(i * 0.45),
            Inches(5.0), Inches(0.35),
            f"✓  {item}", sz=11, color=DARK_TEXT)

    # 进行中
    _rrect(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(4.5),
           fill=WHITE, border=WARN)
    _rect(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(0.5), fill=WARN)
    _tb(slide, Inches(7.1), Inches(2.17), Inches(4), Inches(0.4),
        "○  进行中 / 待完成", sz=13, color=WHITE, bold=True)

    todo_items = [
        "DDL 后端 API 对接",
        "AI 接入通义千问 LLM",
        "深色模式主题切换",
        "成绩详情独立页面",
        "校历页面",
        "通知系统",
        "数据库完整迁移",
        "用户记忆 & 个性化推荐",
    ]
    for i, item in enumerate(todo_items):
        _tb(slide, Inches(7.3), Inches(2.75) + Inches(i * 0.45),
            Inches(5.0), Inches(0.35),
            f"○  {item}", sz=11, color=DARK_TEXT)


# ============================================================
# S17 — 总结与展望
# ============================================================
def s17_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page_bg(slide)
    _page_title(slide, "项目亮点与下一步计划")
    _page_number(slide, 17, TOTAL_SLIDES)

    # 左：项目亮点
    _tb(slide, Inches(0.7), Inches(1.2), Inches(3), Inches(0.4),
        "★  项目亮点", sz=15, color=GOLD, bold=True)

    highlights = [
        ("真实数据", "直连教务系统", TENC_BLUE),
        ("AI 原生", "Function Calling 驱动", PURPLE),
        ("现代全栈", "Next.js + FastAPI + Docker", SUCCESS),
        ("体验优先", "深色主题 + 响应式", WARN),
        ("工程化", "CI/CD + 自动测试", SCU_RED),
    ]
    for i, (title, desc, color) in enumerate(highlights):
        y = Inches(1.8) + Inches(i * 0.85)
        _rrect(slide, Inches(0.7), y, Inches(5.8), Inches(0.65),
               fill=WHITE, border=color, border_w=Pt(0.5))
        _rect(slide, Inches(0.7), y, Inches(0.08), Inches(0.65), fill=color)
        _tb(slide, Inches(1.0), y + Inches(0.12), Inches(1.8), Inches(0.35),
            title, sz=12, color=color, bold=True)
        _tb(slide, Inches(2.8), y + Inches(0.14), Inches(3.5), Inches(0.35),
            desc, sz=11, color=SUB_TEXT)

    # 右：下一步计划
    _tb(slide, Inches(7.0), Inches(1.2), Inches(3), Inches(0.4),
        "▸  下一步计划", sz=15, color=TENC_BLUE, bold=True)

    plans = [
        ("1. AI 对话全面接入", "通义千问 + Function Calling 上线"),
        ("2. DDL 系统完善", "后端 API + 数据库 + 提醒通知"),
        ("3. 个性化记忆系统", "口味偏好、课程习惯、智能推荐"),
        ("4. RAG 文档问答", "校规、选课手册知识库检索"),
        ("5. 全面测试与上线", "性能优化、安全加固、用户测试"),
    ]
    for i, (title, desc) in enumerate(plans):
        y = Inches(1.8) + Inches(i * 0.85)
        _rrect(slide, Inches(7.0), y, Inches(5.8), Inches(0.65),
               fill=RGBColor(0xF0, 0xF5, 0xFF), border=TENC_BLUE, border_w=Pt(0.5))
        _tb(slide, Inches(7.3), y + Inches(0.05), Inches(5.2), Inches(0.3),
            title, sz=12, color=TENC_BLUE, bold=True)
        _tb(slide, Inches(7.3), y + Inches(0.35), Inches(5.2), Inches(0.25),
            desc, sz=10, color=SUB_TEXT)


# ============================================================
# S18 — 致谢页
# ============================================================
def s18_thanks():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=TENC_DARK)
    _rect(slide, 0, 0, Inches(6.5), SH, fill=TENC_BLUE)
    _circle(slide, Inches(10), Inches(-0.5), Inches(4), RGBColor(0x1A, 0x3D, 0x7A))

    _tb(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
        "谢谢！", sz=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
        "SCU Assistant — 让校园生活更智能",
        sz=22, color=RGBColor(0xBB, 0xD5, 0xFF), align=PP_ALIGN.CENTER)
    _rect(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.03), fill=TENC_LIGHT)
    _tb(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.4),
        "四川大学计算机学院  |  2026 春季  |  软件工程课程设计",
        sz=13, color=MID_GRAY, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
        "欢迎提问与交流", sz=14, color=RGBColor(0x88, 0xAA, 0xDD),
        align=PP_ALIGN.CENTER)


# ============================================================
# 生成
# ============================================================
print("生成 PPT 中...")

s01_cover()
s02_toc()
s03_sec_bg()       # Section 1
s04_background()
s05_data()
s06_sec_pos()      # Section 2
s07_positioning()
s08_sec_feat()     # Section 3
s09_features()
s10_sec_arch()     # Section 4
s11_architecture()
s12_techstack()
s13_team()
s14_timeline()
s15_risks()
s16_progress()
s17_summary()
s18_thanks()

out = os.path.join(os.path.dirname(__file__), "SCU_Assistant_开题报告_v2.pptx")
prs.save(out)
print(f"✓ 已生成: {out}")
print(f"  共 {len(prs.slides)} 页")
