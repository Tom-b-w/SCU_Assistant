"""
将 Gemini 生成的 ISPA 研究图转为 PPT
包含两个主要 panel: (a) Label Shift 问题 (b) ISPA 方法
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 配色
# ============================================================
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
SUB_TEXT   = RGBColor(0x64, 0x74, 0x8B)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)

# Panel (a) colors
IDEAL_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
IDEAL_BG     = RGBColor(0xE8, 0xF5, 0xE9)
REAL_ORANGE  = RGBColor(0xE6, 0x5C, 0x00)
REAL_BG      = RGBColor(0xFF, 0xF3, 0xE0)
DANGER_RED   = RGBColor(0xC6, 0x28, 0x28)

# Panel (b) colors
BLUE_MAIN    = RGBColor(0x15, 0x65, 0xC0)
BLUE_LIGHT   = RGBColor(0xE3, 0xF2, 0xFD)
BLUE_BG      = RGBColor(0xBB, 0xDE, 0xFB)
TEAL         = RGBColor(0x00, 0x79, 0x6B)
TEAL_BG      = RGBColor(0xE0, 0xF2, 0xF1)
PURPLE       = RGBColor(0x6A, 0x1B, 0x9A)
PURPLE_BG    = RGBColor(0xF3, 0xE5, 0xF5)
AMBER        = RGBColor(0xF5, 0x7F, 0x17)
GREEN_OK     = RGBColor(0x2E, 0x7D, 0x32)
GREEN_BG     = RGBColor(0xC8, 0xE6, 0xC9)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ============================================================
# 工具函数
# ============================================================
def _rect(slide, left, top, w, h, fill=None, border=None, border_w=Pt(0.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    return s

def _rrect(slide, left, top, w, h, fill=None, border=None, border_w=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    return s

def _circle(slide, left, top, d, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def _tb(slide, left, top, w, h, text, sz=14, color=DARK_TEXT,
        bold=False, italic=False, font="Segoe UI",
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.15):
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    p.line_spacing = Pt(sz * line_spacing)
    tf.auto_size = None
    return tf

def _multi_tb(slide, left, top, w, h, lines, default_sz=12, default_color=DARK_TEXT,
              align=PP_ALIGN.LEFT, font="Segoe UI"):
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(lines):
        text = item[0]
        sz = item[1] if len(item) > 1 else default_sz
        color = item[2] if len(item) > 2 else default_color
        bold = item[3] if len(item) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.line_spacing = Pt(sz * 1.4)
    return tf

def _arrow_right(slide, x, y, w, color=MID_GRAY, h=Inches(0.04)):
    """水平箭头线"""
    _rect(slide, x, y, w, h, fill=color)
    # 箭头头部
    s = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x + w - Inches(0.15), y - Inches(0.08), Inches(0.2), Inches(0.2))
    s.rotation = 90
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def _arrow_down(slide, x, y, h, color=MID_GRAY, w=Inches(0.04)):
    """垂直箭头线"""
    _rect(slide, x, y, w, h, fill=color)


# ============================================================
# Slide 1 — 封面
# ============================================================
def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=RGBColor(0x0D, 0x47, 0xA1))
    _circle(slide, Inches(9), Inches(-1), Inches(5), RGBColor(0x15, 0x65, 0xC0))
    _circle(slide, Inches(10), Inches(4), Inches(6), RGBColor(0x0A, 0x3A, 0x80))

    _tb(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
        "ISPA", sz=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
        "Instance-adaptive Scaling for Prior Adjustment",
        sz=24, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
    _rect(slide, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.03), fill=BLUE_BG)
    _tb(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
        "Label Shift Correction: From Ideal to Real-World Deployment",
        sz=16, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2 — Panel (a): Problem Statement
# ============================================================
def slide_problem():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    # 顶栏
    _rect(slide, 0, 0, SW, Inches(0.85), fill=RGBColor(0x0D, 0x47, 0xA1))
    _tb(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
        "(a) Label Shift Correction: Ideal Assumption vs. Real-World Deployment",
        sz=20, color=WHITE, bold=True)

    # ---- 左半：Ideal Setting ----
    left_x = Inches(0.5)
    _rrect(slide, left_x, Inches(1.2), Inches(5.8), Inches(5.5),
           fill=IDEAL_BG, border=IDEAL_GREEN, border_w=Pt(1.5))
    _tb(slide, left_x + Inches(0.3), Inches(1.35), Inches(5.2), Inches(0.5),
        "Ideal Setting (prior work)  \u2714", sz=18, color=IDEAL_GREEN, bold=True,
        align=PP_ALIGN.CENTER)

    # 公式
    _tb(slide, left_x + Inches(0.5), Inches(2.0), Inches(4.8), Inches(0.5),
        "p(x|y) = p_S(x|y) = p_T(x|y)", sz=16, color=DARK_TEXT, bold=True,
        align=PP_ALIGN.CENTER, font="Consolas")

    # 描述卡片
    _rrect(slide, left_x + Inches(0.4), Inches(2.7), Inches(4.8), Inches(1.5),
           fill=WHITE, border=IDEAL_GREEN)
    _multi_tb(slide, left_x + Inches(0.7), Inches(2.85), Inches(4.2), Inches(1.3), [
        ("Key Assumption", 14, IDEAL_GREEN, True),
        ("Same appearance, different proportions", 12, DARK_TEXT, False),
        ("\u2192 Standard EM works perfectly", 12, IDEAL_GREEN, True),
        ("Global \u03b1 = 1.0 \u2192 ALL predictions CORRECT", 12, IDEAL_GREEN, False),
    ])

    # Source / Target 示意
    _rrect(slide, left_x + Inches(0.4), Inches(4.5), Inches(2.1), Inches(1.8),
           fill=WHITE, border=LIGHT_GRAY)
    _tb(slide, left_x + Inches(0.5), Inches(4.6), Inches(1.9), Inches(0.3),
        "Source Distribution", sz=10, color=SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)
    # 模拟柱状图
    bars_s = [(0.3, BLUE_MAIN), (0.5, BLUE_MAIN), (0.2, BLUE_MAIN)]
    for i, (h, c) in enumerate(bars_s):
        bx = left_x + Inches(0.8) + Inches(i * 0.4)
        bh = Inches(h * 2)
        by = Inches(6.0) - bh
        _rect(slide, bx, by, Inches(0.25), bh, fill=c)
    _tb(slide, left_x + Inches(0.5), Inches(6.05), Inches(1.9), Inches(0.2),
        "0.%   0.%   0.%", sz=7, color=SUB_TEXT, align=PP_ALIGN.CENTER)

    _rrect(slide, left_x + Inches(2.8), Inches(4.5), Inches(2.1), Inches(1.8),
           fill=WHITE, border=LIGHT_GRAY)
    _tb(slide, left_x + Inches(2.9), Inches(4.6), Inches(1.9), Inches(0.3),
        "Target Distribution", sz=10, color=SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)
    bars_t = [(0.2, TEAL), (0.3, TEAL), (0.5, TEAL)]
    for i, (h, c) in enumerate(bars_t):
        bx = left_x + Inches(3.2) + Inches(i * 0.4)
        bh = Inches(h * 2)
        by = Inches(6.0) - bh
        _rect(slide, bx, by, Inches(0.25), bh, fill=c)

    # ---- 右半：Real-World Deployment ----
    right_x = Inches(6.8)
    _rrect(slide, right_x, Inches(1.2), Inches(5.8), Inches(5.5),
           fill=REAL_BG, border=REAL_ORANGE, border_w=Pt(1.5))
    _tb(slide, right_x + Inches(0.3), Inches(1.35), Inches(5.2), Inches(0.5),
        "Real-World Deployment (our focus)  \u2718", sz=18, color=DANGER_RED, bold=True,
        align=PP_ALIGN.CENTER)

    # 公式
    _tb(slide, right_x + Inches(0.5), Inches(2.0), Inches(4.8), Inches(0.5),
        "p(x|y) \u2260 p_S(x|y) \u2260 p_T(x|y)", sz=16, color=DANGER_RED, bold=True,
        align=PP_ALIGN.CENTER, font="Consolas")

    # 描述
    _rrect(slide, right_x + Inches(0.4), Inches(2.7), Inches(4.8), Inches(1.5),
           fill=WHITE, border=REAL_ORANGE)
    _multi_tb(slide, right_x + Inches(0.7), Inches(2.85), Inches(4.2), Inches(1.3), [
        ("Real-World Challenge", 14, DANGER_RED, True),
        ("Same species, different appearance", 12, DARK_TEXT, False),
        ("\u2192 Global correction is catastrophic", 12, DANGER_RED, True),
        ("Effective \u03b1 window: only 0.04 wide!", 12, REAL_ORANGE, False),
    ])

    # Consequence 卡片
    _rrect(slide, right_x + Inches(0.4), Inches(4.5), Inches(2.1), Inches(1.8),
           fill=WHITE, border=DANGER_RED)
    _tb(slide, right_x + Inches(0.5), Inches(4.6), Inches(1.9), Inches(0.3),
        "Consequence", sz=11, color=DANGER_RED, bold=True, align=PP_ALIGN.CENTER)
    _multi_tb(slide, right_x + Inches(0.5), Inches(4.95), Inches(1.9), Inches(1.0), [
        ("F1 collapses", 13, DANGER_RED, True),
        ("59% \u2192 danger zone", 11, DARK_TEXT, False),
        ("Too aggressive / Too conservative", 10, SUB_TEXT, False),
    ], align=PP_ALIGN.CENTER)

    # Razor edge 卡片
    _rrect(slide, right_x + Inches(2.8), Inches(4.5), Inches(2.1), Inches(1.8),
           fill=WHITE, border=AMBER)
    _tb(slide, right_x + Inches(2.9), Inches(4.6), Inches(1.9), Inches(0.3),
        "Razor Edge", sz=11, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    _multi_tb(slide, right_x + Inches(2.9), Inches(4.95), Inches(1.9), Inches(1.0), [
        ("Effective \u03b1 window", 12, AMBER, True),
        ("only 0.04 wide", 14, DANGER_RED, True),
        ("\u2190 Too aggressive | Too conservative \u2192", 9, SUB_TEXT, False),
    ], align=PP_ALIGN.CENTER)


# ============================================================
# Slide 3 — Panel (b): ISPA Method Overview
# ============================================================
def slide_ispa_method():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _rect(slide, 0, 0, SW, Inches(0.85), fill=TEAL)
    _tb(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
        "(b) ISPA: Instance-adaptive Scaling for Prior Adjustment",
        sz=20, color=WHITE, bold=True)

    # ---- Input ----
    _rrect(slide, Inches(0.4), Inches(1.5), Inches(1.5), Inches(2.0),
           fill=BLUE_LIGHT, border=BLUE_MAIN)
    _tb(slide, Inches(0.5), Inches(1.6), Inches(1.3), Inches(0.4),
        "Input", sz=14, color=BLUE_MAIN, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0.5), Inches(2.1), Inches(1.3), Inches(0.5),
        "[Image]", sz=11, color=SUB_TEXT, align=PP_ALIGN.CENTER)
    # 模拟图片占位
    _rrect(slide, Inches(0.7), Inches(2.2), Inches(0.9), Inches(0.9),
           fill=RGBColor(0xC8, 0xE6, 0xC9), border=LIGHT_GRAY)
    _tb(slide, Inches(0.7), Inches(2.4), Inches(0.9), Inches(0.5),
        "\U0001f33f", sz=24, color=GREEN_OK, align=PP_ALIGN.CENTER)

    # Arrow → Frozen Classifier
    _arrow_right(slide, Inches(1.95), Inches(2.5), Inches(0.5), color=MID_GRAY)

    # ---- Frozen Classifier ----
    _rrect(slide, Inches(2.6), Inches(1.5), Inches(1.8), Inches(2.0),
           fill=PURPLE_BG, border=PURPLE)
    _tb(slide, Inches(2.7), Inches(1.6), Inches(1.6), Inches(0.35),
        "Frozen Classifier", sz=12, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(2.7), Inches(2.0), Inches(1.6), Inches(0.3),
        "(ResNet50)", sz=11, color=PURPLE, align=PP_ALIGN.CENTER)
    # 模拟网络图
    _rrect(slide, Inches(3.1), Inches(2.4), Inches(0.8), Inches(0.8),
           fill=WHITE, border=PURPLE)
    _tb(slide, Inches(3.1), Inches(2.5), Inches(0.8), Inches(0.6),
        "\U0001f9e0", sz=22, color=PURPLE, align=PP_ALIGN.CENTER)

    # Arrow →
    _arrow_right(slide, Inches(4.5), Inches(2.5), Inches(0.5), color=MID_GRAY)

    # ---- Stage 1: Prior Estimation ----
    s1_x = Inches(5.1)
    _rrect(slide, s1_x, Inches(1.2), Inches(3.0), Inches(2.8),
           fill=BLUE_LIGHT, border=BLUE_MAIN, border_w=Pt(1.5))
    _rect(slide, s1_x, Inches(1.2), Inches(3.0), Inches(0.45), fill=BLUE_MAIN)
    _tb(slide, s1_x + Inches(0.1), Inches(1.25), Inches(2.8), Inches(0.4),
        "Stage 1: Prior Estimation", sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _rrect(slide, s1_x + Inches(0.2), Inches(1.85), Inches(2.5), Inches(0.5),
           fill=WHITE, border=BLUE_MAIN)
    _tb(slide, s1_x + Inches(0.3), Inches(1.9), Inches(2.3), Inches(0.4),
        "Regularized EM", sz=12, color=BLUE_MAIN, bold=True, align=PP_ALIGN.CENTER)

    _multi_tb(slide, s1_x + Inches(0.3), Inches(2.5), Inches(2.5), Inches(1.2), [
        ("w_k: target estimation", 11, DARK_TEXT, False),
        ("w_k = \u03c0_k / c_k", 12, BLUE_MAIN, True),
        ("\u03c0\u0302 = \u03bb\u00b7q_EM + (1-\u03bb)\u00b7u", 11, DARK_TEXT, False),
    ], font="Consolas")

    # Arrow →
    _arrow_right(slide, Inches(8.2), Inches(2.5), Inches(0.5), color=MID_GRAY)

    # ---- Stage 2: Confidence-Adaptive Weights ----
    s2_x = Inches(8.8)
    _rrect(slide, s2_x, Inches(1.2), Inches(4.0), Inches(2.8),
           fill=TEAL_BG, border=TEAL, border_w=Pt(1.5))
    _rect(slide, s2_x, Inches(1.2), Inches(4.0), Inches(0.45), fill=TEAL)
    _tb(slide, s2_x + Inches(0.1), Inches(1.25), Inches(3.8), Inches(0.4),
        "Stage 2: Confidence-Adaptive Weights", sz=13, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)

    _tb(slide, s2_x + Inches(0.2), Inches(1.8), Inches(3.6), Inches(0.35),
        "Confidence Gauge", sz=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # 三个置信度示例
    conf_data = [
        ("95%", "High confidence", GREEN_OK, GREEN_BG),
        ("15%", "Low confidence", AMBER, RGBColor(0xFF, 0xF8, 0xE1)),
        ("45%", "Medium confidence", BLUE_MAIN, BLUE_LIGHT),
    ]
    for i, (pct, label, color, bg) in enumerate(conf_data):
        cy = Inches(2.2) + Inches(i * 0.55)
        _rrect(slide, s2_x + Inches(0.3), cy, Inches(3.3), Inches(0.45),
               fill=bg, border=color)
        _tb(slide, s2_x + Inches(0.5), cy + Inches(0.05), Inches(0.6), Inches(0.35),
            pct, sz=12, color=color, bold=True)
        # 进度条
        bar_w = Inches(1.5)
        _rect(slide, s2_x + Inches(1.2), cy + Inches(0.12), bar_w, Inches(0.2),
              fill=LIGHT_GRAY)
        fill_pct = int(pct.replace("%", "")) / 100
        _rect(slide, s2_x + Inches(1.2), cy + Inches(0.12), int(bar_w * fill_pct), Inches(0.2),
              fill=color)
        _tb(slide, s2_x + Inches(2.8), cy + Inches(0.05), Inches(0.8), Inches(0.35),
            label, sz=8, color=SUB_TEXT)

    # ---- 底部：ISPA 核心公式与可视化 ----
    # 公式区域
    _rrect(slide, Inches(0.4), Inches(4.3), Inches(6.5), Inches(2.8),
           fill=RGBColor(0xF5, 0xF5, 0xF5), border=BLUE_MAIN)
    _rect(slide, Inches(0.4), Inches(4.3), Inches(6.5), Inches(0.45), fill=BLUE_MAIN)
    _tb(slide, Inches(0.6), Inches(4.35), Inches(6.1), Inches(0.4),
        "ISPA Core Formula", sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _multi_tb(slide, Inches(0.8), Inches(4.9), Inches(5.8), Inches(1.8), [
        ("\u03c6 = 0:  No correction applied  (\u03b1_i \u2248 1)", 12, DARK_TEXT, False),
        ("\u03c6 = 1:  Full correction  (\u03b1_i \u2248 s)", 12, DARK_TEXT, False),
        ("", 8, DARK_TEXT, False),
        ("ISPA \u03b1_i: Instance-adaptive", 14, BLUE_MAIN, True),
        ("Strength varied per sample based on confidence \u03c6_i", 12, TEAL, False),
        ("", 8, DARK_TEXT, False),
        ("\u03c6 = (1-conf)^b  \u2192  High conf \u2192 less correction", 11, PURPLE, False),
    ], font="Consolas")

    # 可视化区域
    _rrect(slide, Inches(7.2), Inches(4.3), Inches(5.5), Inches(2.8),
           fill=RGBColor(0xF5, 0xF5, 0xF5), border=TEAL)
    _rect(slide, Inches(7.2), Inches(4.3), Inches(5.5), Inches(0.45), fill=TEAL)
    _tb(slide, Inches(7.4), Inches(4.35), Inches(5.1), Inches(0.4),
        "ISPA Correctness Visualization (showing \u03b1_i adaptation)", sz=13,
        color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 四个 alpha 状态可视化
    alpha_states = [
        ("\u03b1_i \u2248 s", "Confidently\nCORRECT", GREEN_OK, GREEN_BG, "ADAPTED &\nCORRECT"),
        ("\u03b1_i \u2248 s", "Adapted", TEAL, TEAL_BG, "Moderate\nconfidence"),
        ("\u03b1_i \u2248 1", "Preserved", AMBER, RGBColor(0xFF, 0xF8, 0xE1), "Low\nconfidence"),
        ("\u03b1_i \u2248 1", "No change", MID_GRAY, RGBColor(0xF0, 0xF0, 0xF0), "Uncertain"),
    ]
    for i, (alpha, label, color, bg, desc) in enumerate(alpha_states):
        ax = Inches(7.5) + Inches(i * 1.3)
        ay = Inches(4.95)
        _rrect(slide, ax, ay, Inches(1.1), Inches(1.8), fill=bg, border=color)
        _tb(slide, ax + Inches(0.05), ay + Inches(0.1), Inches(1.0), Inches(0.3),
            alpha, sz=10, color=color, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
        _tb(slide, ax + Inches(0.05), ay + Inches(0.5), Inches(1.0), Inches(0.5),
            label, sz=9, color=color, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, ax + Inches(0.05), ay + Inches(1.1), Inches(1.0), Inches(0.6),
            desc, sz=8, color=SUB_TEXT, align=PP_ALIGN.CENTER)

    # 底部 aggregated correct count
    _rrect(slide, Inches(7.5), Inches(6.85), Inches(5.0), Inches(0.4),
           fill=GREEN_BG, border=GREEN_OK)
    _tb(slide, Inches(7.6), Inches(6.88), Inches(4.8), Inches(0.35),
        "\u2192 Aggregated correct count increases with ISPA", sz=11,
        color=GREEN_OK, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 4 — 关键优势总结
# ============================================================
def slide_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill=WHITE)
    _rect(slide, 0, 0, SW, Inches(0.85), fill=RGBColor(0x0D, 0x47, 0xA1))
    _tb(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
        "ISPA: Key Advantages", sz=20, color=WHITE, bold=True)

    advantages = [
        ("Instance-Adaptive",
         "Per-sample \u03b1_i instead of global \u03b1",
         "Avoids the catastrophic failure of one-size-fits-all correction",
         BLUE_MAIN, BLUE_LIGHT),
        ("Confidence-Driven",
         "Confidence gauge \u03c6 controls correction strength",
         "High-confidence samples get full correction; uncertain ones stay safe",
         TEAL, TEAL_BG),
        ("Regularized EM",
         "Smoothed prior estimation with uniform fallback",
         "Prevents extreme weight collapse from noisy or biased class estimates",
         PURPLE, PURPLE_BG),
        ("Robust to Appearance Shift",
         "Works when p_S(x|y) \u2260 p_T(x|y)",
         "The real-world scenario where standard methods catastrophically fail",
         REAL_ORANGE, REAL_BG),
    ]

    for i, (title, subtitle, desc, color, bg) in enumerate(advantages):
        y = Inches(1.3) + Inches(i * 1.4)
        _rrect(slide, Inches(0.8), y, Inches(11.5), Inches(1.2),
               fill=bg, border=color, border_w=Pt(1))
        _rect(slide, Inches(0.8), y, Inches(0.1), Inches(1.2), fill=color)
        _tb(slide, Inches(1.2), y + Inches(0.1), Inches(4), Inches(0.4),
            title, sz=16, color=color, bold=True)
        _tb(slide, Inches(1.2), y + Inches(0.5), Inches(5), Inches(0.3),
            subtitle, sz=12, color=DARK_TEXT, font="Consolas")
        _tb(slide, Inches(6.5), y + Inches(0.3), Inches(5.5), Inches(0.5),
            desc, sz=12, color=SUB_TEXT)

    # 底部核心信息
    _rrect(slide, Inches(2), Inches(6.3), Inches(9), Inches(0.7),
           fill=RGBColor(0x0D, 0x47, 0xA1))
    _tb(slide, Inches(2.2), Inches(6.4), Inches(8.6), Inches(0.5),
        "ISPA: Don't correct what you're not confident about.",
        sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# 生成
# ============================================================
print("生成 ISPA PPT 中...")

slide_cover()
slide_problem()
slide_ispa_method()
slide_summary()

out = os.path.join(os.path.dirname(__file__), "ISPA_Method_Overview.pptx")
prs.save(out)
print(f"\u2713 已生成: {out}")
print(f"  共 {len(prs.slides)} 页")
