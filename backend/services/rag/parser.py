"""文档解析：PDF/PPT/TXT -> 文本块"""
import io
import logging
import zipfile
from xml.etree import ElementTree

logger = logging.getLogger(__name__)


def chunk_text(text: str, max_chars: int = 800, overlap: int = 100) -> list[str]:
    """将长文本按段落/句子边界切分为重叠块。"""
    text = text.strip()
    if not text:
        return []
    if len(text) <= max_chars:
        return [text]

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = start + max_chars
        if end >= len(text):
            chunks.append(text[start:].strip())
            break
        # 尝试在句号/换行处断句
        best = -1
        for sep in ["\n\n", "\n", "。", ". ", "；", "！", "？"]:
            pos = text.rfind(sep, start + max_chars // 2, end)
            if pos > best:
                best = pos + len(sep)
        if best <= start:
            best = end
        chunks.append(text[start:best].strip())
        start = best - overlap
    return [c for c in chunks if c]


def parse_pdf_bytes(data: bytes, filename: str = "") -> str:
    """解析 PDF 文件字节流，返回全文文本。"""
    import fitz  # PyMuPDF
    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n\n".join(pages)


def parse_pptx_bytes(data: bytes, filename: str = "") -> str:
    """解析 PPTX 文件字节流，返回全文文本。"""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip("| "):
                        parts.append(row_text)
        slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


def parse_docx_bytes(data: bytes, filename: str = "") -> str:
    """解析 DOCX 文档字节流，返回正文和表格文本。"""
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        xml_data = archive.read("word/document.xml")

    root = ElementTree.fromstring(xml_data)
    parts: list[str] = []

    for paragraph in root.findall(".//w:p", namespace):
        texts = [
            node.text
            for node in paragraph.findall(".//w:t", namespace)
            if node.text
        ]
        line = "".join(texts).strip()
        if line:
            parts.append(line)

    return "\n".join(parts)


def parse_file(data: bytes, filename: str) -> str:
    """根据文件扩展名选择解析器。"""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pdf":
        return parse_pdf_bytes(data, filename)
    elif ext in ("pptx", "ppt"):
        return parse_pptx_bytes(data, filename)
    elif ext == "docx":
        return parse_docx_bytes(data, filename)
    elif ext in ("txt", "md"):
        return data.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"不支持的文件类型: .{ext}")
